
# -*- coding: utf-8 -*-
"""
Статистическая лаборатория v5.6

Это приложение предназначено для проведения разведочного анализа данных (EDA),
выполнения простых A/B‑тестов и построения базовых прогнозов по временным
рядам. В версии 5.6 добавлено множество улучшений по качеству данных,
поддержке большего числа форматов, улучшены отчёты и расширены
возможности по автоматическому подбору параметров.

Ключевые изменения v5.6:
* Автоматическое приведение типов при загрузке файлов (попытка распознать
  числа и даты в строковых колонках), поддержка форматов CSV, Excel,
  Parquet и Feather.
* Улучшенный генератор Excel‑отчётов: добавлен титульный лист (Cover),
  проверка типов входных аргументов и устойчивость к передаче `None` или
  произвольных объектов вместо DataFrame.
* Импутация и очистка: возможность заполнять пропуски медианой или
  наиболее частым значением, а также выполнять винзоризацию на уровне
  столбцов.
* Расширенный A/B‑модуль: поддержка Mann–Whitney U и z‑теста для
  пропорций, вывод размера эффекта, бизнес‑сводка.
* Модуль временных рядов: добавлена автоматическая настройка ARIMA
  посредством перебора параметров (auto ARIMA) и вывод метрик
  кросс‑валидации (MAE, MAPE, RMSE).
* Быстрые пресеты для выбора колонок: можно мгновенно выбрать все
  числовые колонки без идентификаторов, либо топ‑10 наиболее
  волатильных метрик по коэффициенту вариации.
* Один клик — один отчёт: возможность скачать zip‑архив с HTML‑ и
  Excel‑отчётами, а также конфиг аналитики (filters + выбранные колонки).
Запуск: streamlit run stats_lab_app_v5_5.py
"""
import streamlit as st
import warnings
import pandas as pd
import numpy as np
import json
from io import BytesIO

# pptx для генерации отчётов PowerPoint
try:
    from pptx import Presentation  # type: ignore
    from pptx.util import Inches, Pt  # type: ignore
except Exception:
    Presentation = None  # если библиотека недоступна, отчёты PPTX не будут созданы
import math
# Внешние библиотеки: scipy, plotly, statsmodels
# Их можно установить: pip install scipy plotly statsmodels
from scipy import stats
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
try:
    from statsmodels.tsa.seasonal import STL
except Exception:
    STL = None  # аккуратно обработаем ниже

try:
    from statsmodels.tsa.arima.model import ARIMA
except Exception:
    ARIMA = None  # используем только если доступно

# Пытаемся определить наличие pmdarima для продвинутого авто-ARIMA. Если нет,
# переменная будет False. Но авто‑ARIMA реализовано собственными средствами,
# поэтому доступность pmdarima не критична.
try:
    import pmdarima  # type: ignore
    PMDARIMA_AVAILABLE = True
except Exception:
    PMDARIMA_AVAILABLE = False

# --------------
# ------------------------- UI STYLE -------------------------
def inject_custom_css():
    st.markdown(
        """
        <style>
        .business-summary {
            border-radius: 12px;
            border: 1px solid #d1fae5;
            background-color: #ecfdf5;
            padding: .8rem 1rem;
            font-size: .95rem;
            white-space: pre-line;
        }

        .stTabs [data-baseweb="tab"] {
            border-radius: 10px;
        }
        .stTabs [aria-selected="true"] {
            background: #eef4ff !important;
        }

        .step-badge {
            padding: 0.35rem 0.7rem;
            border-radius: 999px;
            background-color: #eef2ff;
            color: #1f2933;
            font-size: 0.8rem;
            display: inline-flex;
            align-items: center;
            gap: 0.35rem;
        }
        .step-badge span {
            font-weight: 600;
        }
        .step-badge .muted {
            opacity: .8;
            font-weight: 400;
        }
        </style>
        """,
        unsafe_allow_html=True,
    )


METHOD_INFO = {
    "mean": {"name": "Среднее", "description": "Арифметическое среднее.", "when": "Симметричные распределения."},
    "median": {"name": "Медиана", "description": "Середина упорядоченных наблюдений.", "when": "Есть выбросы/скос."},
    "mode": {"name": "Мода", "description": "Наиболее частое значение.", "when": "Категориальные/дискретные данные."},
    "std": {"name": "Std", "description": "Стандартное отклонение.", "when": "Оценка вариативности."},
    "cv": {"name": "CV", "description": "Относительный разброс (std/mean).", "when": "Сравнение стабильности."},
    "iqr": {"name": "IQR", "description": "Межквартильный размах.", "when": "Устойчив к выбросам."},
}


# ------------------------- HELPERS -------------------------
@st.cache_data
def describe_basic_stats(series: pd.Series):
    col = series.dropna()
    res = {
        "count": int(col.count()),
        "mean": float(col.mean()) if col.count() else np.nan,
        "median": float(col.median()) if col.count() else np.nan,
        "min": float(col.min()) if col.count() else np.nan,
        "max": float(col.max()) if col.count() else np.nan,
        "range": float(col.max() - col.min()) if col.count() else np.nan,
        "var": float(col.var(ddof=1)) if col.count() > 1 else np.nan,
        "std": float(col.std(ddof=1)) if col.count() > 1 else np.nan,
        "q1": float(col.quantile(0.25)) if col.count() else np.nan,
        "q3": float(col.quantile(0.75)) if col.count() else np.nan,
    }
    res["iqr"] = res["q3"] - res["q1"] if not np.isnan(res["q3"]) and not np.isnan(res["q1"]) else np.nan
    res["cv"] = (res["std"] / res["mean"]) if res["mean"] not in (0, np.nan) and not np.isnan(res["mean"]) else np.nan

    try:
        mode_values = stats.mode(col, keepdims=True)
        res["mode"] = float(mode_values.mode[0]) if len(mode_values.mode) > 0 else np.nan
    except Exception:
        res["mode"] = np.nan

    res["skewness"] = float(stats.skew(col, bias=False)) if len(col) > 2 else np.nan
    res["kurtosis"] = float(stats.kurtosis(col, fisher=True, bias=False)) if len(col) > 3 else np.nan
    for p in (5, 10, 25, 50, 75, 90, 95):
        res[f"p{p}"] = float(col.quantile(p/100)) if len(col) else np.nan
    return res


def detect_outliers_iqr(series, k=1.5):
    col = series.dropna()
    if len(col) == 0:
        return pd.Series(False, index=series.index), np.nan, np.nan, np.nan
    q1, q3 = col.quantile(0.25), col.quantile(0.75)
    iqr = q3 - q1
    lower, upper = q1 - k * iqr, q3 + k * iqr
    mask = (series < lower) | (series > upper)
    return mask, float(lower), float(upper), float(iqr)


def detect_outliers_z(series, z_thresh=3.0):
    col = series.dropna()
    if len(col) == 0 or col.std(ddof=0) == 0:
        return pd.Series(False, index=series.index)
    z = (series - col.mean()) / col.std(ddof=0)
    return z.abs() > z_thresh

# ------------------------- IMPUTATION & CLEANING HELPERS -------------------------
def winsorize_series(series: pd.Series, lower_pct: float = 0.01, upper_pct: float = 0.99) -> pd.Series:
    """
    Клиппинг (винзоризация) значений в серии по квантилям.
    Значения ниже lower_pct заменяются на квантиль lower_pct,
    а выше upper_pct — на квантиль upper_pct.

    Параметры:
        series: исходный числовой столбец;
        lower_pct: нижний квантиль (0–1);
        upper_pct: верхний квантиль (0–1).

    Возвращает:
        Series с подрезанными значениями.
    """
    # Если колонка полностью состоит из NaN или нечисловых значений
    if series.dropna().empty:
        return series
    low = series.quantile(lower_pct)
    high = series.quantile(upper_pct)
    return series.clip(lower=low, upper=high)

def impute_dataframe(df: pd.DataFrame, columns: list[str], strategy: str = "median") -> pd.DataFrame:
    """
    Выполняет импутацию пропусков в указанных колонках по выбранной стратегии.

    Поддерживаемые стратегии:
      - 'median': заполнение медианой (для числовых колонок).
      - 'most_frequent': заполнение наиболее частым значением.

    Возвращает новый DataFrame с заполненными значениями (оригинал не изменяется).
    """
    new_df = df.copy()
    for c in columns:
        if strategy == "median":
            try:
                med = new_df[c].median()
                new_df[c] = new_df[c].fillna(med)
            except Exception:
                pass
        elif strategy == "most_frequent":
            try:
                mode = new_df[c].mode()
                if not mode.empty:
                    new_df[c] = new_df[c].fillna(mode.iloc[0])
            except Exception:
                pass
    return new_df

# -----------------------------------------------------------------------------
# Типизация и преобразование
# -----------------------------------------------------------------------------
def auto_coerce_dataframe_types(df: pd.DataFrame, date_threshold: float = 0.8, num_threshold: float = 0.8) -> pd.DataFrame:
    """
    Пытается преобразовать строковые столбцы в числовые или datetime, если
    подавляющее большинство значений можно интерпретировать как такие типы.

    :param df: исходный DataFrame
    :param date_threshold: минимальная доля успешных преобразований, чтобы колонка
        считалась датой
    :param num_threshold: минимальная доля успешных преобразований, чтобы колонка
        считалась числовой
    :return: новый DataFrame с приведёнными типами (оригинал не изменяется)
    """
    new_df = df.copy()
    for col in new_df.columns:
        s = new_df[col]
        # Проверяем только object/string колонки
        if not pd.api.types.is_object_dtype(s):
            continue
        try:
            # Пробуем привести к числу. Если в числе используются запятые как
            # разделитель дробной части, заменяем их на точки.
            str_series = s.astype(str).str.replace(",", ".", regex=False)
            numeric_converted = pd.to_numeric(str_series, errors="coerce")
            ratio_num = numeric_converted.notna().sum() / len(numeric_converted) if len(numeric_converted) > 0 else 0.0
            if ratio_num >= num_threshold:
                new_df[col] = numeric_converted
                continue
            # Пробуем привести к дате (подавляем предупреждения о формате)
            with warnings.catch_warnings():
                warnings.simplefilter("ignore", UserWarning)
                date_converted = pd.to_datetime(s, errors="coerce")
            ratio_date = date_converted.notna().sum() / len(date_converted) if len(date_converted) > 0 else 0.0
            if ratio_date >= date_threshold:
                new_df[col] = date_converted
                continue
        except Exception:
            pass
    return new_df


def normality_test(series):
    col = series.dropna()
    if len(col) < 3:
        return None
    sample = col.sample(5000, random_state=42) if len(col) > 5000 else col
    stat, p_value = stats.shapiro(sample)
    return {"statistic": float(stat), "p_value": float(p_value), "n_used": int(len(sample))}


def cohen_d(x, y):
    x, y = np.array(x), np.array(y)
    nx, ny = len(x), len(y)
    if nx < 2 or ny < 2:
        return np.nan
    sx, sy = x.std(ddof=1), y.std(ddof=1)
    sp = np.sqrt(((nx - 1) * sx**2 + (ny - 1) * sy**2) / (nx + ny - 2))
    return float((x.mean() - y.mean()) / sp) if sp else np.nan


def detect_id_columns(df, threshold_ratio=0.9, min_unique=10):
    n = len(df)
    if n == 0: return []
    out = []
    nunq = df.nunique(dropna=False)
    for c, k in nunq.items():
        if k >= min_unique and k / n >= threshold_ratio:
            out.append(c)
    return out


@st.cache_data
def compute_corr_matrix_cached(df: pd.DataFrame, cols: tuple, method: str = "pearson"):
    return df[list(cols)].corr(method=method)


@st.cache_data
def compute_data_quality_table(df: pd.DataFrame):
    n = len(df)
    rows = []
    for col in df.columns:
        s = df[col]
        n_missing = int(s.isna().sum())
        missing_pct = float((n_missing / n * 100) if n else 0.0)
        n_unique = int(s.nunique(dropna=True))
        is_constant = bool(n_unique <= 1)
        base_dtype = str(s.dtype)
        type_suggestion: str | None = None
        mixed_types = False
        # Анализируем только строковые или object колонки
        if pd.api.types.is_object_dtype(s):
            try:
                # Попытка определить, можно ли привести колонку к числу
                str_series = s.astype(str).str.replace(",", ".", regex=False)
                num_converted = pd.to_numeric(str_series, errors="coerce")
                ratio_num = num_converted.notna().sum() / len(num_converted) if len(num_converted) > 0 else 0.0
                # Попытка привести к дате (подавляем предупреждения о формате)
                with warnings.catch_warnings():
                    warnings.simplefilter("ignore", UserWarning)
                    date_converted = pd.to_datetime(s, errors="coerce")
                ratio_date = date_converted.notna().sum() / len(date_converted) if len(date_converted) > 0 else 0.0
                # Если одна из долей достаточно велика, рекомендуем привести тип
                if ratio_num >= 0.8:
                    type_suggestion = "numeric"
                elif ratio_date >= 0.8:
                    type_suggestion = "datetime"
                # Если часть значений распозналась, но менее 80%, считаем колонку смешанной
                if (0.0 < ratio_num < 0.8) or (0.0 < ratio_date < 0.8):
                    mixed_types = True
            except Exception:
                pass
        rows.append(
            {
                "column": col,
                "dtype": base_dtype,
                "n_missing": n_missing,
                "missing_%": missing_pct,
                "n_unique": n_unique,
                "is_constant": is_constant,
                "mixed_types": mixed_types,
                "type_suggestion": type_suggestion,
            }
        )
    return pd.DataFrame(rows)


def get_strong_correlations(corr_matrix, threshold=0.7):
    rec = []
    cols = corr_matrix.columns
    for i in range(len(cols)):
        for j in range(i+1, len(cols)):
            r = corr_matrix.iloc[i, j]
            if np.isnan(r): continue
            if abs(r) >= threshold:
                rec.append({"feature_1": cols[i], "feature_2": cols[j], "r": float(r), "abs_r": float(abs(r))})
    return pd.DataFrame(rec).sort_values("abs_r", ascending=False) if rec else pd.DataFrame(
        columns=["feature_1", "feature_2", "r", "abs_r"]
    )


def maybe_downsample_xy(x, y, max_points=10000):
    if len(x) <= max_points: return x, y
    idx = np.random.choice(len(x), size=max_points, replace=False)
    return (x.iloc[idx] if isinstance(x, pd.Series) else x[idx],
            y.iloc[idx] if isinstance(y, pd.Series) else y[idx])


@st.cache_data
def compute_acf(series, max_lag):
    x = series.dropna().values
    n = len(x)
    if n == 0:
        return np.arange(max_lag + 1), np.full(max_lag + 1, np.nan)
    x = x - x.mean()
    denom = np.dot(x, x)
    if denom == 0:
        return np.arange(max_lag + 1), np.full(max_lag + 1, np.nan)
    acf_vals = [1.0]
    for lag in range(1, max_lag + 1):
        num = np.dot(x[:-lag], x[lag:])
        acf_vals.append(num / denom)
    return np.arange(max_lag + 1), np.array(acf_vals)


def generate_ts_features(df, date_col, value_col, window=7, spike_thresh_pct=50.0):
    data = df[[date_col, value_col]].dropna().copy().sort_values(date_col)
    data[date_col] = pd.to_datetime(data[date_col])

    data["lag_1"] = data[value_col].shift(1)
    data["diff_1"] = data[value_col] - data["lag_1"]
    data["pct_change_1"] = data[value_col].pct_change() * 100.0
    data["rolling_mean_window"] = data[value_col].rolling(window=window, min_periods=1).mean()
    data["rolling_std_window"] = data[value_col].rolling(window=window, min_periods=1).std(ddof=1)
    data["rolling_cv_window"] = (data["rolling_std_window"] / data["rolling_mean_window"]).replace([np.inf, -np.inf], np.nan)
    data["spike_flag"] = data["pct_change_1"].abs() > spike_thresh_pct

    y = data[value_col].values
    x = np.arange(len(y))
    slope, intercept = (np.polyfit(x, y, 1) if (len(y) >= 2 and len(np.unique(y)) > 1) else (np.nan, np.nan))

    global_features = {
        "n_points": int(len(y)),
        "mean": float(np.nanmean(y)) if len(y) > 0 else np.nan,
        "std": float(np.nanstd(y, ddof=1)) if len(y) > 1 else np.nan,
        "cv": float(np.nanstd(y, ddof=1) / np.nanmean(y)) if len(y) > 1 and np.nanmean(y) != 0 else np.nan,
        "slope_trend": float(slope),
        "intercept_trend": float(intercept),
        "first_value": float(y[0]) if len(y) > 0 else np.nan,
        "last_value": float(y[-1]) if len(y) > 0 else np.nan,
        "change_abs": float(y[-1] - y[0]) if len(y) > 1 else np.nan,
        "change_pct": float((y[-1] - y[0]) / y[0] * 100.0) if len(y) > 1 and y[0] != 0 else np.nan,
    }
    return data, global_features


def plot_ts_plotly(df_ts, date_col, value_col, method, window=7):
    data = df_ts[[date_col, value_col]].dropna().copy().sort_values(date_col)
    data[date_col] = pd.to_datetime(data[date_col])

    plot_df = pd.DataFrame({"date": data[date_col], "value": data[value_col]})
    if method == "Скользящее среднее":
        plot_df["smoothed"] = plot_df["value"].rolling(window=window, min_periods=1).mean()
    elif method == "EWMA":
        plot_df["smoothed"] = plot_df["value"].ewm(span=window, adjust=False).mean()
    elif method == "Скользящая медиана":
        plot_df["smoothed"] = plot_df["value"].rolling(window=window, min_periods=1).median()
    else:
        plot_df["smoothed"] = np.nan

    fig = px.line(plot_df, x="date", y=["value", "smoothed"], labels={"date": "Дата"})
    fig.update_layout(legend_title_text="Серии", hovermode="x unified", margin=dict(t=40, r=20, b=40, l=40))
    return fig


# ------------------------- BUSINESS TEXT -------------------------

def ts_forecast_arima(df, date_col, value_col, horizon, order=(1, 1, 1)):
    """
    Простой прогноз ARIMA для одномерного временного ряда.
    Возвращает исторические данные и прогноз с доверительным интервалом.
    """
    if ARIMA is None:
        raise ImportError("ARIMA недоступен: установите statsmodels>=0.12.")

    data = df[[date_col, value_col]].dropna().copy()
    if data.empty:
        raise ValueError("Недостаточно данных для прогноза.")

    data[date_col] = pd.to_datetime(data[date_col])
    data = data.sort_values(date_col).set_index(date_col)
    y = pd.to_numeric(data[value_col], errors="coerce").dropna()
    if len(y) < 5:
        raise ValueError("Недостаточно точек (минимум 5) для построения ARIMA-прогноза.")

    model = ARIMA(y, order=order)
    res = model.fit()

    forecast_res = res.get_forecast(steps=int(horizon))
    mean_forecast = forecast_res.predicted_mean
    conf_int = forecast_res.conf_int(alpha=0.1)  # 90% ДИ

    fc_df = pd.DataFrame(
        {
            "date": mean_forecast.index.to_timestamp() if hasattr(mean_forecast.index, "to_timestamp") else mean_forecast.index,
            "forecast": mean_forecast.values,
            "lower": conf_int.iloc[:, 0].values,
            "upper": conf_int.iloc[:, 1].values,
        }
    )

    hist_df = y.reset_index().rename(columns={value_col: "value", date_col: "date"})
    return hist_df, fc_df

# -----------------------------------------------------------------------------
# Auto-ARIMA: перебор параметров для выбора лучшей модели
# -----------------------------------------------------------------------------
def ts_forecast_auto_arima(
    df: pd.DataFrame,
    date_col: str,
    value_col: str,
    horizon: int,
    max_p: int = 2,
    max_d: int = 1,
    max_q: int = 2,
) -> tuple[pd.DataFrame, pd.DataFrame, tuple[int, int, int]]:
    """
    Автоматический подбор параметров ARIMA по минимальному значению BIC.
    Перебирает порядок (p,d,q) в заданных диапазонах и выбирает комбинацию
    с наименьшим значением BIC на тренировочных данных. Если найдено
    несколько комбинаций, выбирается первая. Возвращает исторические
    значения, прогноз и выбранный порядок.

    :param df: датафрейм с временным рядом
    :param date_col: имя столбца дат
    :param value_col: имя столбца значений
    :param horizon: количество точек прогноза
    :param max_p: максимальное значение p
    :param max_d: максимальное значение d
    :param max_q: максимальное значение q
    :returns: (hist_df, forecast_df, best_order)
    """
    if ARIMA is None:
        raise ImportError("ARIMA недоступен: установите statsmodels>=0.12.")
    # Подготовка данных
    data = df[[date_col, value_col]].dropna().copy()
    if data.empty:
        raise ValueError("Недостаточно данных для прогноза.")
    data[date_col] = pd.to_datetime(data[date_col])
    data = data.sort_values(date_col).set_index(date_col)
    y = pd.to_numeric(data[value_col], errors="coerce").dropna()
    if len(y) < 10:
        raise ValueError("Недостаточно точек (минимум 10) для авто-ARIMA.")
    # Разделяем данные на обучение и тест (хотя BIC оценивается на всём наборе)
    best_bic = np.inf
    best_order: tuple[int, int, int] = (1, 1, 1)
    best_res = None
    # Ограничиваем d, т.к. перебор d>2 редко осмыслен
    for p in range(max_p + 1):
        for d in range(max_d + 1):
            for q in range(max_q + 1):
                # Пропускаем комбинацию (0,0,0)
                if p == 0 and d == 0 and q == 0:
                    continue
                try:
                    model = ARIMA(y, order=(p, d, q))
                    res = model.fit()
                    bic = res.bic if hasattr(res, "bic") else np.inf
                    if not np.isnan(bic) and bic < best_bic:
                        best_bic = bic
                        best_order = (p, d, q)
                        best_res = res
                except Exception:
                    # пропускаем неудачные комбинации
                    continue
    # Если ничего не подошло (best_res остается None) — используем (1,1,1)
    if best_res is None:
        best_order = (1, 1, 1)
        best_res = ARIMA(y, order=best_order).fit()
    # Прогноз
    forecast_res = best_res.get_forecast(steps=int(horizon))
    mean_forecast = forecast_res.predicted_mean
    conf_int = forecast_res.conf_int(alpha=0.1)
    fc_df = pd.DataFrame(
        {
            "date": mean_forecast.index.to_timestamp() if hasattr(mean_forecast.index, "to_timestamp") else mean_forecast.index,
            "forecast": mean_forecast.values,
            "lower": conf_int.iloc[:, 0].values,
            "upper": conf_int.iloc[:, 1].values,
        }
    )
    hist_df = y.reset_index().rename(columns={value_col: "value", date_col: "date"})
    return hist_df, fc_df, best_order


def business_summary_for_series(col_name, stats_dict, norm_res, n_outliers):
    mean_ = stats_dict.get("mean")
    median_ = stats_dict.get("median")
    cv = stats_dict.get("cv")
    skew = stats_dict.get("skewness")
    txt = [f"По метрике **{col_name}**:"]

    if not np.isnan(mean_) and not np.isnan(median_):
        txt.append("- среднее и медиана **близки**." if abs(mean_ - median_) / (abs(median_) + 1e-9) <= 0.2
                   else "- среднее и медиана **заметно отличаются** → влияет скос/выбросы.")
    if not np.isnan(cv):
        txt.append("- метрика **очень стабильна**." if cv < 0.1 else
                   "- метрика **умеренно изменчива**." if cv < 0.3 else
                   "- **высокая волатильность**, возможны скачки.")
    if not np.isnan(skew):
        txt.append("- хвост вправо (редкие большие значения)." if skew > 0.5 else
                   "- хвост влево (иногда очень низкие значения)." if skew < -0.5 else
                   "- перекоса по хвостам не выявлено.")
    if n_outliers > 0:
        txt.append(f"- обнаружены выбросы: ~{n_outliers} записей.")
    if norm_res and "p_value" in norm_res:
        txt.append("- распределение **нормально** (по Шапиро)." if norm_res["p_value"] >= 0.05
                   else "- распределение **не нормально** → используйте медиану/IQR.")
    return "\n".join(txt)


def business_summary_for_correlation(col1, col2, r, p):
    if np.isnan(r): return "Связь не оценена (мало данных)."
    av = abs(r)
    strength = "почти отсутствует" if av < 0.1 else "слабая" if av < 0.3 else "умеренная" if av < 0.7 else "сильная"
    sign = "прямая" if r > 0 else "обратная"
    out = f"Между **{col1}** и **{col2}** — **{strength} {sign} связь** (r≈{r:.2f})."
    out += " Статистически подтверждена." if p < 0.05 else " Статистически слабая."
    return out


def business_summary_for_ts(global_feats):
    txt = []
    slope = global_feats.get("slope_trend"); change_pct = global_feats.get("change_pct"); cv = global_feats.get("cv")
    if not np.isnan(slope):
        txt.append("- тренд растущий." if slope > 0 else "- тренд нисходящий." if slope < 0 else "- тренд не выражен.")
    if not np.isnan(change_pct):
        sign = "+" if change_pct >= 0 else ""
        txt.append(f"- изменение от старта до конца: **{sign}{change_pct:.1f}%**.")
    if not np.isnan(cv):
        txt.append("- ряд стабильный." if cv < 0.1 else "- умеренные колебания." if cv < 0.3 else "- ряд волатилен.")
    return "\n".join(txt)


def business_summary_for_ab(group_a, group_b, mean_a, mean_b, diff, p_val, alpha, d_value):
    txt = [
        f"Mean **{group_a}** ≈ {mean_a:.2f}, Mean **{group_b}** ≈ {mean_b:.2f}.",
        f"Разница (B - A) ≈ {diff:.2f}."
    ]
    txt.append("Различия **значимы**." if p_val < alpha else "Значимых различий нет.")
    if not np.isnan(d_value):
        ad = abs(d_value)
        eff = "очень маленький" if ad < 0.2 else "небольшой" if ad < 0.5 else "средний" if ad < 0.8 else "большой"
        txt.append(f"Размер эффекта (d≈{d_value:.2f}) → {eff}.")
    return "\n".join(txt)


def build_excel_report(
    df: pd.DataFrame,
    stats_df: pd.DataFrame | None = None,
    corr_matrix: pd.DataFrame | None = None,
    dq: pd.DataFrame | None = None,
    text_blocks: dict | None = None,
) -> BytesIO:
    """Собирает расширенный Excel-отчёт:
    - лист Data: сами данные;
    - лист BasicStats: описательная статистика;
    - лист Correlations: корреляции;
    - лист DataQuality: качество данных;
    - лист AI_Summary: текстовые сводки (EDA/бизнес).
    """
    output = BytesIO()
    with pd.ExcelWriter(output, engine="openpyxl") as writer:
        # Титульный лист (Cover)
        try:
            cover_rows: list[dict] = []
            # Основные размеры
            cover_rows.append({"Параметр": "Строк", "Значение": int(df.shape[0])})
            cover_rows.append({"Параметр": "Столбцов", "Значение": int(df.shape[1])})
            # Резюме из текстовых блоков (обрезаем длинные строки)
            if text_blocks:
                for name, txt in text_blocks.items():
                    if isinstance(txt, str) and txt.strip():
                        cover_rows.append({"Параметр": f"{name}", "Значение": txt[:200] + ("…" if len(txt) > 200 else "")})
            cover_df = pd.DataFrame(cover_rows)
            cover_df.to_excel(writer, sheet_name="Cover", index=False)
        except Exception:
            # даже если Cover создать не удалось, продолжаем
            pass
        # Данные
        try:
            df.to_excel(writer, sheet_name="Data", index=False)
        except Exception:
            # На случай, если df не является DataFrame
            if isinstance(df, pd.DataFrame):
                df.to_excel(writer, sheet_name="Data", index=False)
        # BasicStats
        if isinstance(stats_df, pd.DataFrame) and not stats_df.empty:
            try:
                stats_df.to_excel(writer, sheet_name="BasicStats")
            except Exception:
                pass
        # Correlations
        if isinstance(corr_matrix, pd.DataFrame) and not corr_matrix.empty:
            try:
                corr_matrix.to_excel(writer, sheet_name="Correlations")
            except Exception:
                pass
        # DataQuality
        if isinstance(dq, pd.DataFrame) and not dq.empty:
            try:
                dq.to_excel(writer, sheet_name="DataQuality", index=False)
            except Exception:
                pass
        # AI_Summary
        if text_blocks:
            try:
                rows = [
                    {"section": k, "text": v}
                    for k, v in text_blocks.items()
                    if isinstance(v, str) and v.strip()
                ]
                if rows:
                    summary_df = pd.DataFrame(rows)
                    summary_df.to_excel(writer, sheet_name="AI_Summary", index=False)
            except Exception:
                pass
    output.seek(0)
    return output

def build_pptx_report(
    df: pd.DataFrame,
    stats_df: pd.DataFrame | None = None,
    corr_matrix: pd.DataFrame | None = None,
    dq: pd.DataFrame | None = None,
    summary_text: str = "",
) -> BytesIO:
    """
    Создаёт простой PPTX-отчёт. Содержит титульный слайд, слайд со сводкой,
    таблицы статистики и корреляций, а также рекомендации по качеству данных.
    Возвращает BytesIO с содержимым файла. Если библиотека pptx недоступна —
    возвращает пустой файл.
    """
    # Если библиотека pptx не установлена — возвращаем пустой буфер
    if Presentation is None:
        return BytesIO()
    prs = Presentation()
    # Титульный слайд
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Отчёт: Авто-EDA"
    subtitle = slide.placeholders[1]
    subtitle.text = f"Всего строк: {df.shape[0]}, столбцов: {df.shape[1]}"

    # Слайд сводки
    if summary_text:
        layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = "Сводка"
        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        for line in summary_text.split("\n"):
            p = body.add_paragraph()
            p.text = line
            p.font.size = Pt(12)

    # Слайд базовых статистик
    if isinstance(stats_df, pd.DataFrame) and not stats_df.empty:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = "Базовые статистики"
        # Ограничиваем количество строк/столбцов для читаемости
        table_rows = min(len(stats_df.index) + 1, 15)
        table_cols = min(len(stats_df.columns) + 1, 8)
        display_df = stats_df.copy()
        # Форматируем числовые значения
        for c in display_df.columns:
            display_df[c] = display_df[c].apply(
                lambda x: f"{x:.3g}" if isinstance(x, (int, float, np.floating)) and not pd.isna(x) else str(x)
            )
        display_df = display_df.iloc[: table_rows - 1, : table_cols - 1]
        tbl = slide.shapes.add_table(
            rows=table_rows,
            cols=table_cols,
            left=Inches(0.5),
            top=Inches(2),
            width=Inches(9),
            height=Inches(4),
        ).table
        # Заголовки
        tbl.cell(0, 0).text = "Метрика"
        for j, col in enumerate(display_df.columns):
            tbl.cell(0, j + 1).text = str(col)
        # Заполняем строки
        for i, idx in enumerate(display_df.index):
            tbl.cell(i + 1, 0).text = str(idx)
            for j, col in enumerate(display_df.columns):
                tbl.cell(i + 1, j + 1).text = str(display_df.loc[idx, col])

    # Слайд корреляций
    if isinstance(corr_matrix, pd.DataFrame) and not corr_matrix.empty:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = "Корреляции"
        strong = get_strong_correlations(corr_matrix, threshold=0.5)
        if not strong.empty:
            strong = strong.head(10)[["feature_1", "feature_2", "r"]]
            rows = len(strong.index) + 1
            cols = 3
            tbl = slide.shapes.add_table(
                rows=rows,
                cols=cols,
                left=Inches(0.5),
                top=Inches(2),
                width=Inches(9),
                height=Inches(4),
            ).table
            tbl.cell(0, 0).text = "X"
            tbl.cell(0, 1).text = "Y"
            tbl.cell(0, 2).text = "r"
            for i, (_, row) in enumerate(strong.iterrows()):
                tbl.cell(i + 1, 0).text = str(row["feature_1"])
                tbl.cell(i + 1, 1).text = str(row["feature_2"])
                tbl.cell(i + 1, 2).text = f"{row['r']:.2f}"
        else:
            body = slide.shapes.placeholders[1].text_frame
            body.clear()
            p = body.add_paragraph()
            p.text = "Сильные корреляции не обнаружены."
            p.font.size = Pt(12)

    # Слайд качества данных
    if isinstance(dq, pd.DataFrame) and not dq.empty:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = "Качество данных"
        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        hi = dq[dq["missing_%"] > 30]
        if not hi.empty:
            p = body.add_paragraph()
            p.text = "Колонки с пропусками >30%:"
            p.font.bold = True
            for _, row in hi.iterrows():
                p = body.add_paragraph()
                p.text = f"• {row['column']} ({row['missing_%']:.1f}%)"
        else:
            p = body.add_paragraph()
            p.text = "Пропусков >30% не найдено."
        suggestions = dq[(dq["type_suggestion"].notna()) | (dq["mixed_types"] == True)]
        if not suggestions.empty:
            p = body.add_paragraph()
            p.text = "Рекомендации по типам:"
            p.font.bold = True
            for _, row in suggestions.iterrows():
                if row["type_suggestion"]:
                    p = body.add_paragraph()
                    p.text = f"• {row['column']}: привести к {row['type_suggestion']} (текущий тип {row['dtype']})"
                if row["mixed_types"]:
                    p = body.add_paragraph()
                    p.text = f"• {row['column']}: смешанные типы значений, желательно очистить формат"
    # Сохраняем
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer
def generate_ai_text_with_fallback(base_text: str, extra_prompt: str, api_key: str | None) -> str:
    """
    Пробует дополнить локальную сводку с помощью внешней LLM (OpenAI).
    Работает и с новой (1.x), и со старой (0.x) версией библиотеки openai.
    Если ключ не задан или что-то пошло не так, просто возвращает исходный текст.
    """
    if not api_key:
        # AI-режим не включён — работаем только на локальных эвристиках
        return base_text

    try:
        import openai  # type: ignore

        prompt = extra_prompt + "\n\n" + base_text

        # Новый клиент (openai>=1.x)
        if hasattr(openai, "OpenAI"):
            client = openai.OpenAI(api_key=api_key)
            resp = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[
                    {"role": "system", "content": "Ты аналитик данных. Пиши кратко, по делу, на русском."},
                    {"role": "user", "content": prompt},
                ],
                max_tokens=600,
            )
            content = resp.choices[0].message.content

        # Старый формат (openai 0.x)
        else:
            openai.api_key = api_key
            resp = openai.ChatCompletion.create(
                model="gpt-4o-mini",
                messages=[
                    {"role": "system", "content": "Ты аналитик данных. Пиши кратко, по делу, на русском."},
                    {"role": "user", "content": prompt},
                ],
                max_tokens=600,
            )
            content = resp["choices"][0]["message"]["content"]

        return content.strip() if content else base_text

    except Exception:
        # Ничего не ломаем, просто тихо остаёмся на локальной сводке
        return base_text

def get_ai_config():
    mode = st.session_state.get("ai_mode", "local")
    api_key = st.session_state.get("ai_api_key") if mode == "openai" else None
    return mode, api_key



def render_ai_block(local_text: str, button_label: str, cache_key: str, extra_prompt: str = ""):
    """
    Универсальный блок: локальное резюме + кнопка вызова AI + вывод AI-резюме.
    cache_key – уникальный ключ для session_state (например 'ai_dq', 'ai_ts', 'ai_ab_main').
    """
    if not local_text:
        return

    st.markdown("🔎 **Локальное резюме:**")
    st.markdown(f'<div class="business-summary">{local_text}</div>', unsafe_allow_html=True)

    mode, api_key = get_ai_config()
    if mode != "openai" or not api_key:
        st.caption("AI-резюме доступно в режиме OpenAI на шаге 3.")
        return

    if "ai_summaries" not in st.session_state:
        st.session_state["ai_summaries"] = {}

    if st.button(button_label, key=f"btn_{cache_key}"):
        ai_text = generate_ai_text_with_fallback(local_text, extra_prompt, api_key)
        st.session_state["ai_summaries"][cache_key] = ai_text

    if cache_key in st.session_state["ai_summaries"]:
        st.markdown("🤖 **AI-резюме:**")
        st.markdown(
            f'<div class="business-summary">{st.session_state["ai_summaries"][cache_key]}</div>',
            unsafe_allow_html=True,
        )


def ai_enrich_text(base_text: str, extra_prompt: str = "") -> str:
    """
    УСТАРЕВШАЯ обёртка для AI: сохранена для обратной совместимости.
    Больше НЕ выполняет никаких запросов к внешним моделям и всегда
    возвращает исходный текст.

    Вся работа с AI теперь идёт через функцию render_ai_block, которая
    вызывается только по нажатию соответствующей кнопки в интерфейсе.
    """
    return base_text


def auto_eda_summary(df: pd.DataFrame, stats_df: pd.DataFrame, corr: pd.DataFrame | None, dq: pd.DataFrame, cols: list[str]) -> str:
    """Генерирует текстовую авто-сводку по выбранным числовым колонкам.
    Работает на простых эвристиках, без внешних моделей."""
    txt: list[str] = []
    if not cols:
        return "Нет выбранных колонок для EDA."

    # 1) Самые нестабильные и самые стабильные метрики по CV
    if "cv" in stats_df.columns:
        cv_sorted = stats_df["cv"].abs().sort_values().dropna()
        if not cv_sorted.empty:
            low_name = cv_sorted.index[0]
            low_val = cv_sorted.iloc[0]
            high_name = cv_sorted.index[-1]
            high_val = cv_sorted.iloc[-1]
            txt.append(f"- Самая стабильная метрика: **{low_name}** (CV≈{low_val:.2f}).")
            if len(cv_sorted) > 1:
                txt.append(f"- Наиболее волатильная метрика: **{high_name}** (CV≈{high_val:.2f}).")

    # 2) Метрики с наибольшими пропусками
    if "missing_%" in dq.columns:
        hi_missing = dq.sort_values("missing_%", ascending=False).head(3)
        serious = hi_missing[hi_missing["missing_%"] > 20]
        if not serious.empty:
            probs = ", ".join(f"{idx} ({row['missing_%']:.1f}%)" for idx, row in serious.iterrows())
            txt.append(f"- Колонки с заметной долей пропусков: {probs}.")
        elif (dq["missing_%"] > 0).any():
            txt.append("- Пропуски есть, но их доля во всех колонках < 20%.")
        else:
            txt.append("- Пропуски практически отсутствуют во всех выбранных колонках.")

    # 3) Сильные корреляции между метриками
    if corr is not None and not corr.empty:
        cm = corr.copy()
        for c in cm.columns:
            cm.loc[c, c] = 0.0
        strong_pairs = []
        for i in cm.columns:
            for j in cm.columns:
                if j <= i:
                    continue
                r = cm.loc[i, j]
                if abs(r) >= 0.5:
                    strong_pairs.append((i, j, r))
        if strong_pairs:
            strong_pairs.sort(key=lambda x: -abs(x[2]))
            top_desc = ", ".join(f"{a}–{b} (r≈{r:.2f})" for a, b, r in strong_pairs[:3])
            txt.append(f"- Обнаружены сильные связи между метриками: {top_desc}.")
        else:
            txt.append("- Сильных корреляций между выбранными метриками не выявлено.")

    # 4) Скос и тяжёлые хвосты
    skew_cols = []
    heavy_tail = []
    if "skewness" in stats_df.columns:
        for idx, row in stats_df.iterrows():
            skew = row.get("skewness")
            if isinstance(skew, (int, float)):
                if skew > 1:
                    skew_cols.append(f"{idx} (вправо)")
                elif skew < -1:
                    skew_cols.append(f"{idx} (влево)")
    if "kurtosis" in stats_df.columns:
        for idx, row in stats_df.iterrows():
            kurt = row.get("kurtosis")
            if isinstance(kurt, (int, float)) and kurt > 4:
                heavy_tail.append(f"{idx} (kurt≈{kurt:.1f})")
    if skew_cols:
        txt.append("- Метрики со значительным скосом распределения: " + ", ".join(skew_cols) + ".")
    if heavy_tail:
        txt.append("- Метрики с тяжёлыми хвостами (много экстремальных значений): " + ", ".join(heavy_tail) + ".")

    # 5) Диапазоны и возможные выбросы
    ranges = []
    for c in cols:
        if c in stats_df.index:
            row = stats_df.loc[c]
            min_v = row.get("min", float("nan"))
            max_v = row.get("max", float("nan"))
            if not (math.isnan(min_v) or math.isnan(max_v)):
                ranges.append(f"{c}: [{min_v:.3g}; {max_v:.3g}]")
    if ranges:
        txt.append("- Диапазоны значений по выбранным метрикам: " + "; ".join(ranges) + ".")

    if not txt:
        txt.append("По выбранным метрикам не удалось построить содержательную сводку.")
    return "\n".join(txt)


def build_auto_eda_html(df: pd.DataFrame, cols: list[str], stats_df: pd.DataFrame, corr: pd.DataFrame | None,
                        dq: pd.DataFrame, summary_text: str) -> bytes:
    """Строит HTML-отчёт по авто-EDA с простым оформлением."""
    style = """
    <style>
    body { font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', sans-serif; margin: 20px; background-color: #f7f7fb; }
    h1 { color: #111827; }
    h2 { color: #1f2937; border-bottom: 1px solid #e5e7eb; padding-bottom: 4px; }
    .summary-box { background-color: #eef2ff; border-radius: 8px; padding: 12px 16px; margin-bottom: 16px; }
    .summary-box pre { white-space: pre-wrap; font-family: inherit; margin: 0; }
    table { border-collapse: collapse; margin-top: 8px; }
    th, td { padding: 4px 8px; border: 1px solid #e5e7eb; font-size: 12px; }
    th { background-color: #f3f4f6; }
    </style>
    """
    html_parts: list[str] = []
    html_parts.append("<html><head><meta charset='utf-8'><title>Auto EDA Report</title>" + style + "</head><body>")
    html_parts.append("<h1>Auto EDA отчёт</h1>")
    html_parts.append("<div class='summary-box'><h2>Сводка</h2>")
    html_parts.append("<pre>" + summary_text.replace("<", "&lt;").replace(">", "&gt;") + "</pre></div>")

    html_parts.append("<h2>Базовые статистики</h2>")
    html_parts.append(stats_df.to_html(border=0, classes="table-stats", float_format=lambda x: f"{x:.4g}"))

    if corr is not None and not corr.empty:
        html_parts.append("<h2>Корреляции</h2>")
        html_parts.append(corr.to_html(border=0, classes="table-corr", float_format=lambda x: f"{x:.4g}"))

    html_parts.append("<h2>Качество данных</h2>")
    html_parts.append(dq.to_html(border=0, classes="table-dq", float_format=lambda x: f"{x:.4g}"))

    html_parts.append("</body></html>")
    return "\n".join(html_parts).encode("utf-8")

def make_stl_figure(comp_df: pd.DataFrame):
    fig = make_subplots(rows=4, cols=1, shared_xaxes=True, vertical_spacing=0.02)
    fig.add_trace(go.Scatter(x=comp_df["date"], y=comp_df["observed"], name="Наблюдения"), row=1, col=1)
    fig.add_trace(go.Scatter(x=comp_df["date"], y=comp_df["trend"], name="Тренд"), row=2, col=1)
    fig.add_trace(go.Scatter(x=comp_df["date"], y=comp_df["seasonal"], name="Сезонность"), row=3, col=1)
    fig.add_trace(go.Scatter(x=comp_df["date"], y=comp_df["resid"], name="Остаток"), row=4, col=1)
    fig.update_layout(height=650, title="STL-разложение ряда", hovermode="x unified",
                      margin=dict(t=40, r=20, b=40, l=40))
    return fig


# ------------------------- DEMO DATA -------------------------
def get_demo_dataset(name: str) -> pd.DataFrame:
    """Простые демо-датасеты, чтобы можно было поиграться без своих файлов."""
    np.random.seed(42)
    if name == "Продажи по дням (TS)":
        dates = pd.date_range(end=pd.Timestamp.today().normalize(), periods=120, freq="D")
        base = np.linspace(100, 200, len(dates))
        season = 20 * np.sin(np.linspace(0, 4 * np.pi, len(dates)))
        noise = np.random.normal(0, 10, len(dates))
        sales = base + season + noise
        df = pd.DataFrame({
            "date": dates,
            "sales": sales,
            "channel": np.random.choice(["online", "offline"], size=len(dates)),
        })
        return df
    else:
        # Маркетинговая кампания: A/B, CTR, конверсии
        n = 1000
        df = pd.DataFrame({
            "user_id": np.arange(1, n+1),
            "group": np.random.choice(["A", "B"], size=n),
            "impressions": np.random.randint(10, 200, size=n),
        })
        # Придумаем клики и заказы с небольшим uplift в группе B
        p_click_A, p_click_B = 0.08, 0.11
        p_conv_A, p_conv_B = 0.02, 0.03
        mask_A = df["group"] == "A"
        clicks = np.zeros(n, dtype=int)
        orders = np.zeros(n, dtype=int)
        for i in range(n):
            if mask_A.iloc[i]:
                clicks[i] = np.random.binomial(df["impressions"].iloc[i], p_click_A)
                orders[i] = np.random.binomial(clicks[i], p_conv_A)
            else:
                clicks[i] = np.random.binomial(df["impressions"].iloc[i], p_click_B)
                orders[i] = np.random.binomial(clicks[i], p_conv_B)
        df["clicks"] = clicks
        df["orders"] = orders
        df["revenue"] = df["orders"] * np.random.uniform(50, 150, size=n)
        return df


# ------------------------- APP -------------------------
def main():
    st.set_page_config(page_title="Статистическая лаборатория v5.6", layout="wide")

    if "wizard_step" not in st.session_state:
        st.session_state["wizard_step"] = 1

    # Глобальные настройки AI-помощника
    if "ai_mode" not in st.session_state:
        st.session_state["ai_mode"] = "local"
    if "ai_api_key" not in st.session_state:
        st.session_state["ai_api_key"] = ""
    # Флаг автоматической генерации AI-сводок: выключено по умолчанию
    if "auto_ai_call" not in st.session_state:
        st.session_state["auto_ai_call"] = False

    st.title("📊 Статистическая лаборатория v5.6 (EDA + бизнес-резюме)")

    with st.expander("ℹ️ Как пользоваться приложением", expanded=True):
        st.markdown(
            "Приложение работает в формате **мастера из 4 шагов**:\n"
            "1) **Шаг 1. Данные** — загрузите CSV/Excel или выберите демо-датасет.\n"
            "2) **Шаг 2. Фильтры и конфиг** — задайте фильтры, объедините датасеты и выберите колонки для анализа.\n"
            "3) **Шаг 3. AI-помощник (опционально)** — выберите режим AI и при необходимости введите API-ключ.\n"
            "4) **Шаг 4. Аналитика и отчёты** — используйте вкладки: данные, базовая статистика, корреляции, временные ряды, выбросы, фичи, категории, A/B, сценарии и отчёты."
        )

    # ------------------------- STEP 1: DATA SOURCE -------------------------
    st.sidebar.header("Шаг 1. Источник данных")

    source_mode = st.sidebar.radio(
        "Выберите источник данных:",
        ["Загрузить файл", "Демо-датасет"],
        key="data_source_mode",
    )

    main_file = merge_file = None
    df_raw = df_merge = None

    if source_mode == "Загрузить файл":
        # Поддерживаем основные форматы: CSV, Excel, Parquet, Feather
        main_file = st.sidebar.file_uploader(
            "Основной файл (CSV/XLSX/Parquet/Feather)",
            type=["csv", "xlsx", "xls", "parquet", "feather"],
        )
        merge_file = st.sidebar.file_uploader(
            "Второй файл для merge (опционально)",
            type=["csv", "xlsx", "xls", "parquet", "feather"],
        )

        def read_user_file(file):
            """
            Универсальное чтение пользовательских файлов с автоматической
            обработкой формата и приведение типов.

            Добавлен простой кеш по имени и размеру файла, чтобы не перечитывать
            и не приводить типы на каждом перерендере Streamlit.
            """
            # Простой кеш в session_state: ключ — (name, size)
            cache = st.session_state.setdefault("file_cache", {})
            try:
                file_id = (getattr(file, "name", None), getattr(file, "size", None))
            except Exception:
                file_id = (getattr(file, "name", None), None)
            if file_id in cache:
                return cache[file_id]

            name = file.name.lower()
            df_tmp: pd.DataFrame
            # CSV
            if name.endswith(".csv"):
                # Пробуем интуитивный разбор: если стандартный разбор не удаётся,
                # повторим с python-engine, который умеет сам определять разделитель.
                try:
                    df_tmp = pd.read_csv(file)
                except Exception:
                    try:
                        df_tmp = pd.read_csv(file, engine="python")
                    except Exception as ee:
                        raise ee
            # Parquet
            elif name.endswith(".parquet"):
                df_tmp = pd.read_parquet(file)
            # Feather
            elif name.endswith(".feather"):
                df_tmp = pd.read_feather(file)
            # Excel
            else:
                df_tmp = pd.read_excel(file)
            # Авто-приведение типов
            try:
                df_tmp = auto_coerce_dataframe_types(df_tmp)
            except Exception:
                pass

            cache[file_id] = df_tmp
            st.session_state["file_cache"] = cache
            return df_tmp

        if main_file is not None:
            try:
                df_raw = read_user_file(main_file)
            except Exception as e:
                st.sidebar.error(f"Ошибка чтения основного файла: {e}")

        if merge_file is not None:
            try:
                df_merge = read_user_file(merge_file)
            except Exception as e:
                st.sidebar.error(f"Ошибка чтения второго файла: {e}")
    else:
        st.sidebar.markdown("**Демо-датасет** — можно быстро протестировать возможности лаборатории.")
        demo_name = st.sidebar.selectbox(
            "Выберите демо-набор",
            ["Продажи по дням (TS)", "Маркетинговая кампания (A/B)"],
            key="demo_name_select",
        )
        df_raw = get_demo_dataset(demo_name)
        st.sidebar.success(f"Используется демо-набор: {demo_name}")

    has_data = df_raw is not None

    # ------------------------- WIZARD STEPPER -------------------------
    step = st.session_state.get("wizard_step", 1)
    step = max(1, min(4, int(step)))
    st.session_state["wizard_step"] = step

    st.markdown("### 🚦 Шаги анализа")
    cols = st.columns(4)
    labels = [
        "1. Данные",
        "2. Фильтры и конфиг",
        "3. AI-помощник",
        "4. Аналитика и отчёты",
    ]
    for i, col in enumerate(cols, start=1):
        with col:
            if i < step:
                cls = "step-badge step-done"
                prefix = "✅"
            elif i == step:
                cls = "step-badge step-active"
                prefix = "🟢"
            else:
                cls = "step-badge step-future"
                prefix = "⚪"
            st.markdown(f'<span class="{cls}">{prefix} {labels[i-1]}</span>', unsafe_allow_html=True)

    c_prev, c_next = st.columns([1, 1])
    with c_prev:
        if st.button("⬅️ Назад", disabled=(step <= 1)):
            st.session_state["wizard_step"] = max(1, step - 1)
            # Совместимость с разными версиями Streamlit
            if hasattr(st, "rerun"):
                st.rerun()
            else:
                st.experimental_rerun()
    with c_next:
        if st.button("Вперёд ➡️", disabled=(step >= 4 or not has_data)):
            st.session_state["wizard_step"] = min(4, step + 1)
            # Совместимость с разными версиями Streamlit
            if hasattr(st, "rerun"):
                st.rerun()
            else:
                st.experimental_rerun()

    # Если данных нет — дальше смысла нет идти
    if not has_data:
        st.info("Загрузите файл или выберите демо-датасет на **шаге 1** в сайдбаре, затем нажмите «Вперёд ➡️».")
        return

    # Если мы ещё на шаге 1, но данные уже есть — показываем краткое резюме и просим
    if step == 1:
        st.success(
            f"✅ Данные загружены: **{df_raw.shape[0]}** строк, **{df_raw.shape[1]}** колонок.\n"
            "Перейдите на **Шаг 2 (кнопка «Вперёд ➡️»)**, чтобы настроить фильтры, merge и открыть аналитические вкладки."
        )
        with st.expander("📁 Основной датасет — первые строки", expanded=True):
            st.dataframe(df_raw.head(20))

        if df_merge is not None:
            with st.expander("📁 Второй датасет — первые строки", expanded=True):
                c_m1, c_m2 = st.columns(2)
                c_m1.metric("Строк (2-й файл)", df_merge.shape[0])
                c_m2.metric("Столбцов (2-й файл)", df_merge.shape[1])
                st.dataframe(df_merge.head(20))

        return

    # ------------------------- Шаг 2 и выше: КОНФИГ + ФИЛЬТРЫ + АНАЛИТИКА -------------------------

    # ------------------------- CONFIG + FILTERS -------------------------
    st.sidebar.header("Шаг 2. Конфигурация анализа (опционально)")
    config_file = st.sidebar.file_uploader("Загрузить конфиг фильтров (JSON)", type="json")
    if config_file:
        try:
            loaded = json.load(config_file)
            # Восстанавливаем конфигурацию фильтров
            st.session_state["filter_cols"] = loaded.get("filters", {}).get("filter_cols", [])
            for col, spec in loaded.get("filters", {}).get("per_column", {}).items():
                if spec.get("type") == "numeric":
                    st.session_state[f"filter_mode_{col}"] = spec.get("mode", "Диапазон")
                    if "range" in spec:
                        st.session_state[f"filter_range_{col}"] = tuple(spec["range"])
                    if "values" in spec:
                        st.session_state[f"filter_vals_num_{col}"] = spec["values"]
                else:
                    st.session_state[f"search_{col}"] = spec.get("search", "")
                    st.session_state[f"filter_vals_{col}"] = spec.get("values", [])
            # Восстанавливаем выбранные метрики для отчёта, если были сохранены
            if "selected_report_columns" in loaded:
                st.session_state["report_num_cols"] = loaded["selected_report_columns"]
            st.sidebar.success("Конфиг фильтров загружен.")
        except Exception as e:
            st.sidebar.error(f"Ошибка чтения конфига: {e}")

    # ------------------------- MERGE + ФИЛЬТРЫ (Шаг 2 и 3) -------------------------
    # На шаге 2 пользователь настраивает фильтры и merge.
    # На шаге 3 используем уже подготовленный датасет без повторного UI.

    # Подготовим дефолтные конфиги фильтров
    default_filter_cfg = {"filter_cols": [], "per_column": {}}
    default_filter2_cfg = {"filter_cols": [], "per_column": {}}

    if step == 2:
        # 1) Фильтры для основного датасета
        st.subheader("🔍 Фильтры для основного датасета")
        df_main = df_raw.copy()
        active_filters_main = []

        filter_cols = st.multiselect("Колонки для фильтрации (основной датасет)", df_raw.columns.tolist(),
                                     key="filter_cols")
        filter_config = {"filter_cols": filter_cols, "per_column": {}}

        # 2) Фильтры для второго датасета (если он есть)
        df_second = df_merge.copy() if df_merge is not None else None
        active_filters_second = []

        if df_merge is not None:
            st.subheader("🔍 Фильтры для второго датасета")
            filter2_cols = st.multiselect("Колонки для фильтрации (второй датасет)", df_merge.columns.tolist(),
                                          key="filter2_cols")
        else:
            filter2_cols = []

        # Глобальный сброс фильтров и настроек merge
        if st.button("🔁 Сбросить фильтры и merge"):
            for k in list(st.session_state.keys()):
                if (k.startswith("filter_") or k.startswith("search_") or
                    k.startswith("filter2_") or k.startswith("search2_") or
                    k.startswith("merge_") or k in ("use_merge", "merged_df_cache",
                                                    "df_analysis", "active_filters_all",
                                                    "analysis_cols_default", "filter_config", "filter_config_second")):
                    del st.session_state[k]
            if hasattr(st, "rerun"):
                st.rerun()
            else:
                st.experimental_rerun()

        # Применяем фильтры к основному датасету
        for col in filter_cols:
            s = df_raw[col]
            if np.issubdtype(s.dtype, np.number):
                mode = st.radio(f"Тип фильтра для {col} (основной)", ["Диапазон", "По значениям"],
                                key=f"filter_mode_{col}", horizontal=True)
                cfg = {"type": "numeric", "mode": mode}
                if mode == "Диапазон":
                    min_v, max_v = float(s.min()), float(s.max())
                    if not (np.isnan(min_v) or np.isnan(max_v)):
                        r = st.slider(f"Диапазон для {col} (основной)", min_value=min_v, max_value=max_v,
                                      value=st.session_state.get(f"filter_range_{col}", (min_v, max_v)),
                                      key=f"filter_range_{col}")
                        df_main = df_main[(df_main[col] >= r[0]) & (df_main[col] <= r[1])]
                        active_filters_main.append(f"{col} ∈ [{r[0]:.3g}; {r[1]:.3g}]")
                        cfg["range"] = [r[0], r[1]]
                else:
                    vals = sorted(s.dropna().unique())
                    default = st.session_state.get(f"filter_vals_num_{col}", vals[:1] if vals else [])
                    chosen = st.multiselect(f"Значения для {col} (основной)", options=vals, default=default,
                                            key=f"filter_vals_num_{col}")
                    if chosen:
                        df_main = df_main[df_main[col].isin(chosen)]
                        active_filters_main.append(f"{col} ∈ {{{', '.join(map(str, chosen))}}}")
                        cfg["values"] = list(map(lambda x: x if isinstance(x, (int, float)) else str(x), chosen))
                filter_config["per_column"][col] = cfg
            else:
                uniq = sorted(s.dropna().astype(str).unique())
                search = st.text_input(f"Поиск по {col} (основной)", value=st.session_state.get(f"search_{col}", ""),
                                       key=f"search_{col}")
                choices = [v for v in uniq if search.lower() in v.lower()] if search else uniq
                default = st.session_state.get(f"filter_vals_{col}", choices[:1] if choices else [])
                chosen = st.multiselect(f"Значения для {col} (основной)", options=choices, default=default,
                                        key=f"filter_vals_{col}")
                if chosen:
                    df_main = df_main[df_main[col].astype(str).isin(chosen)]
                    active_filters_main.append(f"{col} ∈ {{{', '.join(chosen)}}}")
                filter_config["per_column"][col] = {"type": "categorical", "search": search, "values": chosen}

        # Применяем фильтры ко второму датасету, если он есть
        filter_config_second = {"filter_cols": [], "per_column": {}}
        if df_second is not None and filter2_cols:
            filter_config_second["filter_cols"] = filter2_cols
            for col in filter2_cols:
                s2 = df_merge[col]
                if np.issubdtype(s2.dtype, np.number):
                    mode2 = st.radio(f"Тип фильтра для {col} (2-й датасет)", ["Диапазон", "По значениям"],
                                     key=f"filter2_mode_{col}", horizontal=True)
                    cfg2 = {"type": "numeric", "mode": mode2}
                    if mode2 == "Диапазон":
                        min2, max2 = float(s2.min()), float(s2.max())
                        if not (np.isnan(min2) or np.isnan(max2)):
                            r2 = st.slider(f"Диапазон для {col} (2-й датасет)", min_value=min2, max_value=max2,
                                           value=st.session_state.get(f"filter2_range_{col}", (min2, max2)),
                                           key=f"filter2_range_{col}")
                            df_second = df_second[(df_second[col] >= r2[0]) & (df_second[col] <= r2[1])]
                            active_filters_second.append(f"[2] {col} ∈ [{r2[0]:.3g}; {r2[1]:.3g}]")
                            cfg2["range"] = [r2[0], r2[1]]
                    else:
                        vals2 = sorted(s2.dropna().unique())
                        default2 = st.session_state.get(f"filter2_vals_num_{col}", vals2[:1] if vals2 else [])
                        chosen2 = st.multiselect(f"Значения для {col} (2-й датасет)", options=vals2, default=default2,
                                                 key=f"filter2_vals_num_{col}")
                        if chosen2:
                            df_second = df_second[df_second[col].isin(chosen2)]
                            active_filters_second.append(f"[2] {col} ∈ {{{', '.join(map(str, chosen2))}}}")
                            cfg2["values"] = list(map(lambda x: x if isinstance(x, (int, float)) else str(x), chosen2))
                    filter_config_second["per_column"][col] = cfg2
                else:
                    uniq2 = sorted(s2.dropna().astype(str).unique())
                    search2 = st.text_input(f"Поиск по {col} (2-й датасет)",
                                            value=st.session_state.get(f"search2_{col}", ""),
                                            key=f"search2_{col}")
                    choices2 = [v for v in uniq2 if search2.lower() in v.lower()] if search2 else uniq2
                    default2 = st.session_state.get(f"filter2_vals_{col}", choices2[:1] if choices2 else [])
                    chosen2 = st.multiselect(f"Значения для {col} (2-й датасет)", options=choices2, default=default2,
                                             key=f"filter2_vals_{col}")
                    if chosen2:
                        df_second = df_second[df_second[col].astype(str).isin(chosen2)]
                        active_filters_second.append(f"[2] {col} ∈ {{{', '.join(chosen2)}}}")
                    filter_config_second["per_column"][col] = {"type": "categorical", "search": search2, "values": chosen2}

        # 3) Merge после фильтрации с кнопкой и join по нескольким колонкам
        st.subheader("🔗 Объединение датасетов (после фильтрации)")
        df = df_main.copy()
        merged_df = None

        if df_second is not None:
            with st.expander("Настройки merge", expanded=False):
                st.caption("Объединение выполняется по уже отфильтрованным датасетам. "
                           "Можно выбрать несколько общих колонок как ключи.")
                common_cols = sorted(set(df_main.columns).intersection(df_second.columns))
                if common_cols:
                    join_keys = st.multiselect("Ключевые колонки (общие для обоих датасетов)",
                                               options=common_cols,
                                               default=st.session_state.get("merge_join_keys", []),
                                               key="merge_join_keys")
                else:
                    join_keys = []
                    st.warning("Нет общих колонок для merge.")

                how_join = st.selectbox("Тип join", ["inner", "left", "right", "outer"], key="merge_how")
                merge_btn = st.button("Выполнить merge отфильтрованных датасетов", key="merge_do")
                use_merge = st.checkbox("Использовать результат merge для дальнейшего анализа",
                                        value=st.session_state.get("use_merge", False), key="use_merge")

                # Выполнение merge только по кнопке
                if merge_btn:
                    if not join_keys:
                        st.error("Выберите хотя бы одну ключевую колонку для merge.")
                    else:
                        try:
                            with st.spinner("Выполняем merge отфильтрованных датасетов..."):
                                merged_tmp = df_main.merge(df_second, on=join_keys, how=how_join)
                            st.session_state["merged_df_cache"] = merged_tmp
                            st.success(f"Успешный merge: {merged_tmp.shape[0]} строк, {merged_tmp.shape[1]} столбцов.")
                            st.dataframe(merged_tmp.head(20))
                        except Exception as e:
                            st.error(f"Ошибка merge: {e}")

                # Если в сессии уже есть результат merge — показываем краткую информацию
                if "merged_df_cache" in st.session_state:
                    merged_cached = st.session_state["merged_df_cache"]
                    st.caption(f"Текущий объединённый датасет: {merged_cached.shape[0]} строк, {merged_cached.shape[1]} столбцов.")
                    with st.expander("Предпросмотр объединённого датасета", expanded=False):
                        st.dataframe(merged_cached.head(20))

        # Выбор датасета для дальнейшего анализа на шаге 2
        if df_second is not None and st.session_state.get("use_merge") and "merged_df_cache" in st.session_state:
            df = st.session_state["merged_df_cache"].copy()
            active_filters = active_filters_main + active_filters_second + [f"MERGE ({st.session_state.get('merge_how', 'inner')})"]
        else:
            df = df_main.copy()
            active_filters = active_filters_main + active_filters_second

        # Считаем "исходные" строки/столбцы для KPI шага 2:
        # если используется merge и есть результат, считаем по объединённому датасету,
        # иначе — по отфильтрованному основному датасету.
        if df_second is not None and st.session_state.get("use_merge") and "merged_df_cache" in st.session_state:
            base_kpi_df = st.session_state["merged_df_cache"]
        else:
            base_kpi_df = df_main

        st.session_state["analysis_kpi_rows_origin"] = int(base_kpi_df.shape[0])
        st.session_state["analysis_kpi_cols_origin"] = int(base_kpi_df.shape[1])

        # Сохраняем выбранный датасет и фильтры в состоянии сессии для использования на шаге 3
        st.session_state["df_analysis"] = df
        st.session_state["active_filters_all"] = active_filters
        st.session_state["filter_config"] = filter_config
        st.session_state["filter_config_second"] = filter_config_second
        st.session_state["analysis_cols_default"] = df.columns.tolist()

    else:
        # Шаг 3 и далее: используем уже подготовленный на шаге 2 датасет
        df = st.session_state.get("df_analysis", df_raw.copy())
        active_filters = st.session_state.get("active_filters_all", [])
        filter_config = st.session_state.get("filter_config", default_filter_cfg)
        filter_config_second = st.session_state.get("filter_config_second", default_filter2_cfg)
    # После применения фильтров формируем списки колонок
    st.sidebar.header("Колонки для анализа")
    # Полный список колонок
    all_columns = df.columns.tolist()
    # Определяем потенциальные ID‑колонки по всему датасету
    id_columns_all = detect_id_columns(df)
    # Все числовые колонки без ID
    numeric_cols_all = [c for c in df.select_dtypes(include=[np.number]).columns if c not in id_columns_all]
    # Рассчитываем коэффициент вариации для каждой числовой колонки (для топ‑10 по CV)
    cv_scores: dict[str, float] = {}
    for c in numeric_cols_all:
        col = df[c].dropna()
        if len(col) == 0:
            cv_scores[c] = 0.0
            continue
        mean_val = col.mean()
        std_val = col.std(ddof=1)
        cv_scores[c] = abs(std_val / mean_val) if mean_val not in (0, np.nan) and not np.isnan(mean_val) else 0.0
    # Блок быстрых пресетов выбора
    st.sidebar.markdown("**Быстрый выбор колонок:**")
    preset_choice = st.sidebar.selectbox(
        "Тип пресета",
        ["—", "Числовые без ID", "Топ‑10 по CV"],
        help="Мгновенный выбор популярных наборов колонок",
        key="preset_choice",
    )
    preset_cols: list[str] | None = None
    if preset_choice == "Числовые без ID":
        preset_cols = numeric_cols_all.copy()
    elif preset_choice == "Топ‑10 по CV":
        # Выбираем до 10 колонок с наибольшим CV
        sorted_cols = sorted(cv_scores.keys(), key=lambda k: cv_scores[k], reverse=True)
        preset_cols = sorted_cols[: min(10, len(sorted_cols))]
    # По умолчанию анализируем все колонки
    analysis_default = preset_cols if preset_cols else all_columns
    analysis_cols = st.sidebar.multiselect(
        "Анализировать колонки",
        options=all_columns,
        default=analysis_default,
    ) or analysis_default

    # Определяем ID‑колонки, числовые и категориальные относительно выбранного набора
    id_columns = detect_id_columns(df[analysis_cols])
    numeric_cols = [c for c in df[analysis_cols].select_dtypes(include=[np.number]).columns if c not in id_columns]
    categorical_cols = [c for c in analysis_cols if c not in numeric_cols]


    # Шаг 3: глобальные настройки AI-помощника (опционально)
    if step == 3:
        st.subheader("Шаг 3. AI-помощник (опционально)")

        mode_label_map = {
            "off": "Отключено",
            "local": "Только локальные эвристики",
            "openai": "OpenAI (через API)",
        }
        reverse_map = {v: k for k, v in mode_label_map.items()}

        current_mode = st.session_state.get("ai_mode", "local")
        ui_label = mode_label_map.get(current_mode, "Только локальные эвристики")

        chosen_label = st.radio(
            "Режим AI-комментариев",
            list(mode_label_map.values()),
            index=list(mode_label_map.values()).index(ui_label),
            horizontal=False,
        )
        st.session_state["ai_mode"] = reverse_map[chosen_label]

        if st.session_state["ai_mode"] == "openai":
            st.session_state["ai_api_key"] = st.text_input(
                "API ключ OpenAI",
                type="password",
                value=st.session_state.get("ai_api_key", ""),
            )
            st.caption("Ключ хранится только в памяти текущей сессии и не сохраняется в файл.")

            if st.button("🔍 Проверить ключ OpenAI"):
                try:
                    import openai  # type: ignore

                    test_prompt = "Ответь одной фразой 'OK', если видишь этот запрос."
                    # Новый клиент (openai>=1.x)
                    if hasattr(openai, "OpenAI"):
                        client = openai.OpenAI(api_key=st.session_state["ai_api_key"])
                        _ = client.chat.completions.create(
                            model="gpt-4o-mini",
                            messages=[
                                {"role": "system", "content": "Проверка доступности модели."},
                                {"role": "user", "content": test_prompt},
                            ],
                            max_tokens=5,
                        )
                    else:
                        openai.api_key = st.session_state["ai_api_key"]
                        _ = openai.ChatCompletion.create(
                            model="gpt-4o-mini",
                            messages=[
                                {"role": "system", "content": "Проверка доступности модели."},
                                {"role": "user", "content": test_prompt},
                            ],
                            max_tokens=5,
                        )
                    st.success("Ключ выглядит рабочим: запрос к OpenAI выполнен успешно.")
                except Exception as e:  # pragma: no cover
                    st.error(f"Не удалось выполнить запрос к OpenAI: {e}")

        # Переключатель автоматической генерации AI-сводок
        st.checkbox(
            "Авто‑генерация AI‑сводок",
            value=st.session_state.get("auto_ai_call", False),
            key="auto_ai_call",
            help=(
                "Когда выключено, AI-сводки не запрашиваются автоматически и отображается только локальная"
                " эвристическая сводка. Включите, если хотите получать AI-комментарии."
            ),
        )

        st.info(
            "Начиная с **шага 4**, все текстовые резюме (EDA, базовая статистика, корреляции, временные ряды, A/B и др.) "
            "могут дополняться комментариями AI при включённой авто‑генерации."
        )
    # Если мы на шаге 2 — показываем только предварительный обзор данных после фильтрации,
    # без перегруженного интерфейса вкладок. Полный анализ доступен на шаге 4.
    if step == 2:
        st.subheader("Шаг 2. Фильтры и конфиг — предварительный просмотр данных")
        c1, c2, c3, c4 = st.columns(4)
        origin_rows = st.session_state.get("analysis_kpi_rows_origin", df_raw.shape[0])
        origin_cols = st.session_state.get("analysis_kpi_cols_origin", df_raw.shape[1])
        c1.metric("Строк (исходно)", origin_rows)
        c2.metric("Строк (фильтр)", df.shape[0])
        c3.metric("Столбцов", df.shape[1])
        c4.metric("Активных фильтров", len(active_filters))

        with st.expander("🔎 Текущие фильтры", expanded=bool(active_filters)):
            if active_filters:
                st.markdown("\n".join(f"- {x}" for x in active_filters))
            else:
                st.info("Фильтры не заданы.")

        with st.expander("Предпросмотр данных после фильтрации", expanded=True):
            st.dataframe(df.head(50))

        if id_columns:
            st.caption(f"ID-подобные колонки (исключены из числовой аналитики): {', '.join(id_columns)}")

        st.info(
            "Когда закончите настройку фильтров, выбора колонок и опциональных AI-настроек, "
            "нажмите кнопку «Вперёд ➡️» в блоке шагов выше, чтобы перейти к полному аналитическому интерфейсу (Шаг 4)."
        )
        return

    # Если мы ещё не дошли до шага 4 — вкладки аналитики не показываем
    if step < 4:
        return

    # ------------------------- TABS -------------------------
    (tab_data, tab_stats, tab_corr, tab_ts, tab_outliers, tab_feats, tab_cats,
     tab_groups, tab_ab, tab_scenarios, tab_dict, tab_report) = st.tabs(
        ["📁 Данные и качество", "📊 Базовая статистика", "🔗 Корреляции",
         "⏱ Временные ряды", "⚠️ Выбросы", "🧩 Фичи", "🏷 Категории",
         "📐 Группы / Pivot", "🧪 A/B", "🔁 Сценарии", "📚 Справочник", "📤 Отчёт"]
    )

    # ---- DATA
    with tab_data:
        st.subheader("📁 Данные (после фильтрации)")
        c1, c2, c3, c4 = st.columns(4)
        origin_rows = st.session_state.get("analysis_kpi_rows_origin", df_raw.shape[0])
        origin_cols = st.session_state.get("analysis_kpi_cols_origin", df_raw.shape[1])
        c1.metric("Строк (исходно)", origin_rows)
        c2.metric("Строк (фильтр)", df.shape[0])
        c3.metric("Столбцов", df.shape[1])
        c4.metric("Активных фильтров", len(active_filters))

        with st.expander("🔎 Текущие фильтры", expanded=bool(active_filters)):
            if active_filters:
                st.markdown("\n".join(f"- {x}" for x in active_filters))
            else:
                st.info("Фильтры не заданы.")

        with st.expander("Таблица данных", expanded=True):
            st.dataframe(df.head(100))

        if id_columns:
            st.caption(f"ID-подобные колонки (исключены из числовой аналитики): {', '.join(id_columns)}")

        st.download_button("📥 Скачать отфильтрованные данные (CSV)",
                           data=df.to_csv(index=False).encode("utf-8"),
                           file_name="filtered_data.csv", mime="text/csv")

        with st.expander("🧼 Качество данных (Data Quality)", expanded=False):
            dq = compute_data_quality_table(df[analysis_cols])
            st.dataframe(dq)
            # Выделяем колонки с высокой долей пропусков
            hi = dq[dq["missing_%"] > 30]
            if hi.empty:
                st.info("Колонок с пропусками более 30% не обнаружено.")
            else:
                st.markdown("**Колонки с пропусками > 30%**")
                st.dataframe(hi)
            # Рекомендации по типам и смешанным значениям
            with st.expander("🛠 Рекомендации по типам", expanded=False):
                recs: list[str] = []
                for _, row in dq.iterrows():
                    if row.get("type_suggestion"):
                        recs.append(
                            f"• {row['column']} — привести к {row['type_suggestion']} (текущий тип {row['dtype']})"
                        )
                    if row.get("mixed_types"):
                        recs.append(
                            f"• {row['column']} — смешанные типы значений, желательно очистить формат"
                        )
                if recs:
                    st.markdown("\n".join(recs))
                else:
                    st.info("Рекомендации по типам отсутствуют.")

            # Краткое локальное резюме по качеству данных + AI-комментарий по кнопке
            issues = []
            many_missing = dq[dq["missing_%"] > 30]["column"].tolist()
            if many_missing:
                issues.append(
                    "Колонки с >30% пропусков: " + ", ".join(many_missing) + " — стоит подумать об исключении или импутации."
                )
            const_cols = dq[dq["is_constant"] == True]["column"].tolist()
            if const_cols:
                issues.append(
                    "Колонки-константы: " + ", ".join(const_cols) + " — практически не несут информации."
                )
            type_suspicious = dq[dq["type_suggestion"].notna()]["column"].tolist()
            if type_suspicious:
                issues.append(
                    "Колонки с подозрительным типом (возможен numeric/datetime): " + ", ".join(type_suspicious) + "."
                )

            local_dq_summary = " ".join(issues) if issues else "Серьёзных проблем с качеством данных не обнаружено."
            render_ai_block(
                local_dq_summary,
                "🤖 AI-комментарий по качеству данных",
                "dq_overall",
                extra_prompt="Сделай краткий комментарий по качеству данных и возможным шагам по очистке."
            )
        st.subheader("💾 Сохранить конфиг анализа")
        # Сохраняем не только фильтры, но и выбранные числовые колонки для отчёта (если уже выбирались)
        cfg_dict = {"filters": filter_config}
        if "report_num_cols" in st.session_state:
            cfg_dict["selected_report_columns"] = st.session_state.get("report_num_cols", [])
        cfg_bytes = json.dumps(cfg_dict, ensure_ascii=False, indent=2).encode("utf-8")
        st.download_button(
            "📥 Скачать конфиг (JSON)",
            data=cfg_bytes,
            file_name="analysis_config.json",
            mime="application/json",
        )

    # ---- BASIC STATS
    with tab_stats:
        st.header("1️⃣ Базовые статистики и распределения")
        if not numeric_cols:
            st.info("Нет числовых столбцов.")
        else:
            col = st.selectbox("Числовая колонка", numeric_cols)
            s = df[col].dropna()
            stats_row = describe_basic_stats(s)

            with st.expander("📊 Таблица статистик", expanded=True):
                st.write(pd.DataFrame(stats_row, index=["значения"]).T)

            norm_res = normality_test(s)
            mask_iqr, *_ = detect_outliers_iqr(s, k=1.5)
            n_out = int(mask_iqr.sum())

            with st.expander("🎯 Доверительный интервал для среднего", expanded=False):
                n = stats_row["count"]
                if n > 1:
                    ci_level = st.slider("Уровень доверия", 0.80, 0.99, 0.95, 0.01)
                    mean_v, std_v = stats_row["mean"], stats_row["std"]
                    se = std_v / np.sqrt(n) if n else np.nan
                    if not np.isnan(se) and se > 0:
                        alpha = 1 - ci_level
                        t_crit = stats.t.ppf(1 - alpha / 2, df=n - 1)
                        ci_low, ci_high = mean_v - t_crit * se, mean_v + t_crit * se
                        st.write({"mean": mean_v, "n": n, "ci_level": ci_level, "ci_low": ci_low, "ci_high": ci_high})
                else:
                    st.info("Недостаточно данных для CI.")

            with st.expander("📈 Гистограмма и распределение", expanded=True):
                if len(s) > 0:
                    fig = px.histogram(s, nbins=30, marginal="box", title=f"Распределение: {col}")
                    fig.update_layout(margin=dict(t=40, r=20, b=40, l=40))
                    st.plotly_chart(fig, use_container_width=True)

            with st.expander("⚖️ Тест Шапиро–Уилка", expanded=False):
                st.write("Недостаточно данных." if norm_res is None else norm_res)

            with st.expander("🧾 Бизнес-резюме по метрике", expanded=True):
                local_summary = business_summary_for_series(col, stats_row, norm_res, n_out)
                render_ai_block(
                    local_summary,
                    "🤖 AI-резюме по метрике",
                    f"metric_{col}",
                    extra_prompt=(
                        f"Сделай краткое бизнес-резюме по метрике '{col}'. "
                        "Опиши её поведение, риски, аномалии и возможные гипотезы для продукта."
                    ),
                )

            with st.expander("📊 EDA по категориям", expanded=False):
                if not categorical_cols:
                    st.info("Нет категориальных колонок.")
                else:
                    gcol = st.selectbox(
                        "Категориальная колонка для сегментации",
                        categorical_cols,
                        key="eda_cat_col",
                    )
                    seg_metrics = st.multiselect(
                        "Числовые метрики по категориям",
                        numeric_cols,
                        default=[col] if col in numeric_cols else numeric_cols[:1],
                        key="eda_cat_metrics",
                    )
                    if not seg_metrics:
                        st.info("Выберите хотя бы одну метрику.")
                    else:
                        agg_df = (
                            df.groupby(gcol)[seg_metrics]
                            .agg(["mean", "median", "count"])
                            .reset_index()
                        )
                        st.dataframe(agg_df.head(50))

                        plot_metric = st.selectbox(
                            "Метрика для графика (среднее по категориям)",
                            seg_metrics,
                            key="eda_cat_plot_metric",
                        )
                        try:
                            mean_by_cat = (
                                df.groupby(gcol)[plot_metric].mean().reset_index()
                            )
                            fig_seg = px.bar(
                                mean_by_cat,
                                x=gcol,
                                y=plot_metric,
                                title=f"Среднее значение {plot_metric} по категориям {gcol}",
                            )
                            st.plotly_chart(fig_seg, use_container_width=True)
                        except Exception as e:
                            st.error(f"Не удалось построить график по категориям: {e}")

                        # Небольшое бизнес-резюме по сегментам
                        top_rows = agg_df.sort_values((seg_metrics[0], "mean"), ascending=False).head(3)
                        bottom_rows = agg_df.sort_values((seg_metrics[0], "mean"), ascending=True).head(3)
                        text_lines = ["EDA по категориям:", f"- Категориальная колонка: {gcol}"]
                        text_lines.append(f"- Метрики: {', '.join(seg_metrics)}")
                        text_lines.append("Топ-3 категории по средней величине первой метрики:")
                        for _, r in top_rows.iterrows():
                            text_lines.append(f"  • {r[gcol]}")
                        text_lines.append("Анти-топ-3 категории:")
                        for _, r in bottom_rows.iterrows():
                            text_lines.append(f"  • {r[gcol]}")
                        seg_summary = "\n".join(text_lines)

                        # Локальное резюме + AI-резюме по кнопке
                        render_ai_block(
                            seg_summary,
                            "🤖 AI-резюме по сегментам",
                            f"segments_{gcol}",
                            extra_prompt=(
                                "Сделай бизнес-комментарий по этим сегментам: "
                                "где метрика выше/ниже, какие гипотезы можно выдвинуть и что протестировать в A/B."
                            ),
                        )



    # ---- CORRELATIONS
    with tab_corr:
        st.header("2️⃣ Корреляционный анализ")
        if numeric_cols and len(numeric_cols) >= 2:
            with st.expander("Матрица корреляций", expanded=True):
                method = st.selectbox("Метод", ["pearson", "spearman"], format_func=lambda m: "Пирсон" if m=="pearson" else "Спирмен")
                cm = compute_corr_matrix_cached(df, tuple(numeric_cols), method)
                st.dataframe(cm)
                fig = px.imshow(cm, text_auto=False, color_continuous_scale="RdBu", zmin=-1, zmax=1,
                                title="Корреляционная матрица")
                st.plotly_chart(fig, use_container_width=True)


            with st.expander("Сильные корреляции", expanded=False):
                thr = st.slider("Порог |r|", 0.0, 1.0, 0.7, 0.05)
                strong = get_strong_correlations(cm, threshold=thr)
                st.write("Нет пар." if strong.empty else strong[["feature_1", "feature_2", "r"]])

            with st.expander("Парная корреляция и scatter", expanded=False):
                c1 = st.selectbox("Первая колонка", numeric_cols, key="corr_c1")
                c2 = st.selectbox("Вторая колонка", numeric_cols, key="corr_c2")
                if c1 != c2:
                    x, y = df[c1].dropna(), df[c2].dropna()
                    idx = x.index.intersection(y.index)
                    x, y = x.loc[idx], y.loc[idx]
                    if len(x) >= 3:
                        pr, pp = stats.pearsonr(x, y)
                        sr, sp = stats.spearmanr(x, y)
                        st.write({"pearson_r": float(pr), "pearson_p": float(pp), "spearman_rho": float(sr), "spearman_p": float(sp)})
                        x1, y1 = maybe_downsample_xy(x, y, 10000)
                        scat = px.scatter(x=x1, y=y1, labels={"x": c1, "y": c2}, title=f"Scatter: {c1} vs {c2}")
                        st.plotly_chart(scat, use_container_width=True)
                        local_corr_summary = business_summary_for_correlation(c1, c2, pr, pp)
                        render_ai_block(
                            local_corr_summary,
                            "🤖 Попросить AI прокомментировать эту связь",
                            f"corr_{c1}_{c2}",
                            extra_prompt=(
                                f"Поясни, что означает корреляция между '{c1}' и '{c2}' для продукта. "
                                "Сделай краткий комментарий для продакт-менеджера."
                            ),
                        )
                else:
                    st.info("Выберите разные колонки.")
        else:
            st.info("Недостаточно числовых колонок.")

    # ---- TIME SERIES
    with tab_ts:
        st.header("3️⃣ Временные ряды")
        date_col = st.selectbox("Дата", ["<нет>"] + analysis_cols)
        if date_col != "<нет>" and numeric_cols:
            value_col = st.selectbox("Метрика", numeric_cols)

            # KPI
            try:
                _, feats = generate_ts_features(df, date_col, value_col, window=7, spike_thresh_pct=50.0)
                c1, c2, c3 = st.columns(3)
                c1.metric("Точек", feats.get("n_points", 0))
                c2.metric("Δ, %", f"{feats.get('change_pct', np.nan):.1f}" if not np.isnan(feats.get("change_pct", np.nan)) else "n/a")
                c3.metric("CV", f"{feats.get('cv', np.nan):.2f}" if not np.isnan(feats.get('cv', np.nan)) else "n/a")
                st.caption("CV — коэффициент вариации (std/mean); помогает сравнивать стабильность разных метрик во времени.")
            except Exception:
                pass

            with st.expander("Сглаживание (Plotly)", expanded=True):
                method = st.selectbox("Метод", ["Скользящее среднее", "EWMA", "Скользящая медиана"])
                window = st.slider("Окно/span", 3, 60, 7, 1)
                show_spikes = st.checkbox("Показывать выбросы (спайки)", value=True)
                spike_thresh = st.slider("Порог изменения, % для спайка", 5.0, 200.0, 50.0, 5.0)
                if st.button("Построить ряд"):
                    try:
                        ts_data, gf = generate_ts_features(
                            df, date_col, value_col, window=window, spike_thresh_pct=spike_thresh
                        )
                        fig_ts = plot_ts_plotly(ts_data, date_col, value_col, method, window)
                        if show_spikes:
                            spikes = ts_data[ts_data["spike_flag"]]
                            if not spikes.empty:
                                fig_ts.add_scatter(
                                    x=spikes[date_col],
                                    y=spikes[value_col],
                                    mode="markers",
                                    name="Спайки",
                                    marker=dict(symbol="circle-open", size=9),
                                )
                        st.plotly_chart(fig_ts, use_container_width=True)
                        with st.expander("Глобальные характеристики", expanded=False):
                            st.json(gf)
                        ts_local_summary = business_summary_for_ts(gf)
                        render_ai_block(
                            ts_local_summary,
                            "🤖 AI-комментарий по временному ряду",
                            f"ts_{value_col}",
                            extra_prompt=(
                                f"Сделай бизнес-резюме по временному ряду для метрики '{value_col}'. "
                                "Опиши тренд, стабильность, изменение во времени и возможные гипотезы."
                            ),
                        )
                    except Exception as e:
                        st.error(f"Ошибка: {e}")


            with st.expander("Прогноз (ARIMA)", expanded=False):
                if ARIMA is None:
                    st.info("ARIMA недоступен: установите statsmodels (tsa.arima).")
                else:
                    # Горизонт прогноза
                    horizon = st.slider("Горизонт прогноза (количество периодов)", 5, 60, 14, 1)
                    # Расширенный список вариантов: добавляем авто-подбор
                    order_label = st.selectbox(
                        "Порядок модели ARIMA (p, d, q)",
                        ["(1, 1, 0)", "(1, 1, 1)", "(2, 1, 1)", "auto (подбор)"],
                    )
                    # Справочник вариантов
                    order_map = {
                        "(1, 1, 0)": (1, 1, 0),
                        "(1, 1, 1)": (1, 1, 1),
                        "(2, 1, 1)": (2, 1, 1),
                    }
                    if st.button("Построить прогноз"):
                        try:
                            # Подбор модели
                            if order_label == "auto (подбор)":
                                # Автоматический перебор параметров ARIMA
                                hist_df, fc_df, auto_order = ts_forecast_auto_arima(
                                    df, date_col, value_col, horizon=int(horizon)
                                )
                                order_desc = f"Авто‑выбранный порядок: {auto_order}"
                            else:
                                # Фиксированный порядок из списка
                                hist_df, fc_df = ts_forecast_arima(
                                    df,
                                    date_col,
                                    value_col,
                                    horizon=int(horizon),
                                    order=order_map[order_label],
                                )
                                auto_order = order_map[order_label]
                                order_desc = f"Использованный порядок: {auto_order}"
                            # Рисуем прогноз
                            fig_f = go.Figure()
                            fig_f.add_trace(
                                go.Scatter(
                                    x=hist_df["date"],
                                    y=hist_df["value"],
                                    mode="lines",
                                    name="История",
                                )
                            )
                            fig_f.add_trace(
                                go.Scatter(
                                    x=fc_df["date"],
                                    y=fc_df["forecast"],
                                    mode="lines",
                                    name="Прогноз",
                                )
                            )
                            fig_f.add_trace(
                                go.Scatter(
                                    x=list(fc_df["date"]) + list(fc_df["date"][::-1]),
                                    y=list(fc_df["upper"]) + list(fc_df["lower"][::-1]),
                                    fill="toself",
                                    mode="lines",
                                    name="ДИ прогноза (90%)",
                                    line=dict(width=0),
                                    opacity=0.3,
                                )
                            )
                            fig_f.update_layout(
                                title=f"Прогноз ARIMA для {value_col}",
                                hovermode="x unified",
                                margin=dict(t=40, r=20, b=40, l=40),
                            )
                            st.plotly_chart(fig_f, use_container_width=True)
                            # Выводим информацию о выбранном порядке
                            st.caption(order_desc)
                            # --- Кросс-валидация и метрики прогноза ---
                            try:
                                # Подготовим данные для CV: сортируем по дате, приводим к числам
                                d_ts = df[[date_col, value_col]].dropna().copy()
                                d_ts[date_col] = pd.to_datetime(d_ts[date_col])
                                d_ts = d_ts.sort_values(date_col)
                                y_all = pd.to_numeric(d_ts[value_col], errors="coerce").dropna()
                                h = int(horizon)
                                # cross-validation только если данных достаточно (≥2*h)
                                if ARIMA is not None and len(y_all) >= 2 * h and h > 0:
                                    train = y_all.iloc[:-h]
                                    test = y_all.iloc[-h:]
                                    # Порядок CV — такой же, как использовался для прогноза
                                    order_cv = auto_order
                                    cv_model = ARIMA(train, order=order_cv)
                                    cv_res = cv_model.fit()
                                    cv_fc = cv_res.forecast(steps=h)
                                    # Метрики
                                    mae = float(np.mean(np.abs(cv_fc.values - test.values)))
                                    # MAPE: избежать деления на ноль
                                    mape_vals = []
                                    for yy_true, yy_pred in zip(test.values, cv_fc.values):
                                        if yy_true != 0:
                                            mape_vals.append(abs((yy_pred - yy_true) / yy_true))
                                    mape = float(np.mean(mape_vals)) if mape_vals else float('nan')
                                    rmse = float(np.sqrt(np.mean((cv_fc.values - test.values)**2)))
                                    st.markdown(
                                        f"**Кросс‑валидация:** MAE≈{mae:.3g}, MAPE≈{mape*100:.2f}%, RMSE≈{rmse:.3g}"
                                    )
                            except Exception:
                                pass

                            # Локальное резюме по прогнозу + AI по кнопке
                            try:
                                last_hist = float(hist_df["value"].iloc[-1])
                                last_fc = float(fc_df["forecast"].iloc[-1])
                                change_abs = last_fc - last_hist
                                change_pct = (change_abs / last_hist * 100.0) if last_hist != 0 else float("nan")
                                fc_lower = float(fc_df["lower"].iloc[-1])
                                fc_upper = float(fc_df["upper"].iloc[-1])
                                direction = "растёт" if change_abs > 0 else "снижается" if change_abs < 0 else "остается на прежнем уровне"
                                parts = [
                                    f"Прогноз для метрики {value_col} на горизонт {horizon} точек: "
                                    f"ключевой показатель {direction} относительно текущего уровня.",
                                    f"Текущее значение: {last_hist:.3g}, прогноз на горизонте: {last_fc:.3g} "
                                    f"(диапазон [{fc_lower:.3g}; {fc_upper:.3g}]).",
                                ]
                                if not np.isnan(change_pct):
                                    sign = "+" if change_pct >= 0 else ""
                                    parts.append(f"Относительное изменение к последнему факту: {sign}{change_pct:.1f}%.")
                                local_fc_summary = " ".join(parts)
                            except Exception:
                                local_fc_summary = f"Прогноз для метрики {value_col} построен."

                            render_ai_block(
                                local_fc_summary,
                                "🤖 AI-резюме по прогнозу",
                                f"ts_forecast_{value_col}",
                                extra_prompt="Интерпретируй этот ARIMA-прогноз для бизнеса: тренд, риски, доверительный интервал и насколько модель надёжна."
                            )
                        except Exception as e:
                            st.error(f"Ошибка прогноза: {e}")

            with st.expander("Ресемплинг", expanded=False):
                freq_label = st.selectbox("Частота", ["D (дни)", "W (недели)", "M (месяцы)"])
                agg_func = st.selectbox("Агрегат", ["mean", "sum", "max", "min"])
                f_map = {"D (дни)": "D", "W (недели)": "W", "M (месяцы)": "M"}
                if st.button("Ресемплировать"):
                    try:
                        d = df[[date_col, value_col]].dropna().copy()
                        d[date_col] = pd.to_datetime(d[date_col])
                        d = d.sort_values(date_col).set_index(date_col)
                        res = getattr(d[value_col].resample(f_map[freq_label]), agg_func)().reset_index()
                        fig = px.line(res, x=date_col, y=value_col, title=f"Ресемплинг ({freq_label}, {agg_func})")
                        st.plotly_chart(fig, use_container_width=True)
                    except Exception as e:
                        st.error(f"Ошибка ресемплинга: {e}")

            with st.expander("STL-разложение", expanded=False):
                if STL is None:
                    st.info("STL недоступен: установите statsmodels.")
                else:
                    period = st.number_input("Период сезонности", 2, 365, 7, 1)
                    if st.button("Выполнить STL"):
                        try:
                            with st.spinner("Считаем STL-декомпозицию..."):
                                d = df[[date_col, value_col]].dropna().copy()
                                d[date_col] = pd.to_datetime(d[date_col])
                                d = d.sort_values(date_col).set_index(date_col)[value_col].asfreq("D").interpolate()
                                res = STL(d, period=int(period), robust=True).fit()
                                comp = pd.DataFrame({"date": d.index, "observed": d.values,
                                                     "trend": res.trend, "seasonal": res.seasonal, "resid": res.resid})
                            st.plotly_chart(make_stl_figure(comp), use_container_width=True)
                        except Exception as e:
                            st.error(f"Ошибка STL: {e}")

            with st.expander("Сезонность по календарю", expanded=False):
                data = df[[date_col, value_col]].dropna().copy()
                if len(data):
                    data[date_col] = pd.to_datetime(data[date_col])
                    gran = st.selectbox("Гранулярность", ["Дни недели", "Месяцы", "Часы"])
                    if gran == "Дни недели":
                        data["key"] = data[date_col].dt.dayofweek
                        name_map = {0:"Пн",1:"Вт",2:"Ср",3:"Чт",4:"Пт",5:"Сб",6:"Вс"}
                    elif gran == "Месяцы":
                        data["key"] = data[date_col].dt.month
                        name_map = {1:"Янв",2:"Фев",3:"Мар",4:"Апр",5:"Май",6:"Июн",7:"Июл",8:"Авг",9:"Сен",10:"Окт",11:"Ноя",12:"Дек"}
                    else:
                        data["key"] = data[date_col].dt.hour
                        name_map = None
                    grp = data.groupby("key")[value_col].mean().rename("value")
                    idx = sorted(grp.index)
                    show = pd.DataFrame({"group": [name_map.get(i, str(i)) if name_map else str(i) for i in idx],
                                         "value": grp.loc[idx].values})
                    fig = px.bar(show, x="group", y="value", title=f"Сезонность ({gran})")
                    st.plotly_chart(fig, use_container_width=True)
                else:
                    st.info("Недостаточно данных.")

            with st.expander("Автокорреляция (ACF)", expanded=False):
                data = df[[date_col, value_col]].dropna().copy()
                if len(data):
                    data = data.sort_values(date_col)
                    s = data[value_col]
                    max_lag = st.slider("Максимальный лаг", 1, max(2, min(60, len(s)-1)), min(20, len(s)-1))
                    if st.button("Построить ACF"):
                        lags, acf_vals = compute_acf(s, max_lag)
                        st.plotly_chart(
                            px.bar(
                                pd.DataFrame({"lag": lags, "acf": acf_vals}),
                                x="lag",
                                y="acf",
                                title=f"ACF для {value_col}",
                            ),
                            use_container_width=True,
                        )
                        # Автоподбор сезонного периода по пику ACF
                        if len(acf_vals) > 2:
                            acf_series = pd.Series(acf_vals, index=lags)
                            acf_series = acf_series[acf_series.index > 1]  # пропускаем лаг 0 и 1
                            if not acf_series.empty:
                                best_lag = int(acf_series.iloc[acf_series.abs().argmax()].name)
                                best_val = float(acf_series.loc[best_lag])
                                if abs(best_val) >= 0.3:
                                    st.success(
                                        f"Возможный сезонный период: **{best_lag}** шагов (ACF≈{best_val:.2f}). "
                                        "Можно использовать это значение как период сезонности в STL или ARIMA."
                                    )
                                else:
                                    st.info("Явно выраженного сезонного периода в ACF не видно.")
                else:
                    st.info("Недостаточно данных для ACF.")

        else:
            if date_col != "<нет>":
                st.info("Нужны числовые столбцы для анализа временных рядов.")

    # ---- OUTLIERS
    with tab_outliers:
        st.header("4️⃣ Поиск выбросов")
        if not numeric_cols:
            st.info("Нет числовых колонок.")
        else:
            out_col = st.selectbox("Колонка", numeric_cols)
            with st.expander("IQR-выбросы", expanded=True):
                if st.button("Найти выбросы (IQR)", help="IQR — межквартильный размах; выбросами считаются точки далеко за пределами [Q1 - 1.5·IQR, Q3 + 1.5·IQR]."):
                    mask, lower, upper, iqr = detect_outliers_iqr(df[out_col], k=1.5)
                    n = int(mask.sum())
                    c1, c2, c3 = st.columns(3)
                    c1.metric("Выбросов", n)
                    c2.metric("% строк", f"{(n/len(df)*100 if len(df) else 0):.2f}%")
                    c3.metric("IQR", f"{iqr:.3g}" if not np.isnan(iqr) else "n/a")
                    if n > 0:
                        out_df = df[mask].copy()
                        with st.expander("Таблица выбросов", expanded=False):
                            st.dataframe(out_df.head(50))
                        st.download_button("📥 Скачать выбросы (CSV)", data=out_df.to_csv(index=False).encode("utf-8"),
                                           file_name=f"outliers_iqr_{out_col}.csv", mime="text/csv")

            with st.expander("Z-score выбросы", expanded=False):
                zt = st.slider(
                    "Порог |Z|",
                    min_value=1.0,
                    max_value=6.0,
                    value=3.0,
                    step=0.5,
                    help="Z‑score — число стандартных отклонений от среднего; используйте порог 3.0 для грубых выбросов и 2.0 для более чувствительного поиска."
                )
                if st.button("Найти выбросы (Z-score)"):
                    mask = detect_outliers_z(df[out_col], z_thresh=zt)
                    n = int(mask.sum())
                    st.metric("Выбросов", n)
                    if n > 0:
                        out_df = df[mask].copy()
                        st.dataframe(out_df.head(50))
                        st.download_button("📥 Скачать выбросы (CSV)", data=out_df.to_csv(index=False).encode("utf-8"),
                                           file_name=f"outliers_z_{out_col}.csv", mime="text/csv")

            # Локальное резюме по выбросам + AI по кнопке
            try:
                mask_iqr, lower_iqr, upper_iqr, iqr_val = detect_outliers_iqr(df[out_col], k=1.5)
                n_iqr = int(mask_iqr.sum())
                mask_z_default = detect_outliers_z(df[out_col], z_thresh=3.0)
                n_z_default = int(mask_z_default.sum())
                parts = [
                    f"Выбросы по колонке {out_col}: IQR‑критерий с k=1.5 даёт {n_iqr} точек "
                    f"(диапазон [{lower_iqr:.3g}; {upper_iqr:.3g}]).",
                    f"Стандартный Z‑score с порогом 3.0 отмечает {n_z_default} наблюдений как выбросы.",
                    "Рекомендуется проверить качество данных, возможные ошибки ввода и, при необходимости, применить клиппинг или импутацию."
                ]
                local_out_summary = " ".join(parts)
            except Exception:
                local_out_summary = f"По колонке {out_col} выбросы рассчитаны, серьёзных ошибок при вычислении не обнаружено."

            render_ai_block(
                local_out_summary,
                "🤖 AI-комментарий по выбросам",
                f"outliers_{out_col}",
                extra_prompt="Прокомментируй масштабы проблемы с выбросами, возможные причины и варианты обработки (клиппинг, фильтрация, импутация)."
            )

            # --- Импутация и очистка данных ---
            with st.expander("Импутация / очистка", expanded=False):
                # выбираем, какие столбцы будем обрабатывать
                cols_imp = st.multiselect("Колонки для обработки", numeric_cols, key="imp_cols")
                method_imp = st.selectbox("Метод", ["median", "most_frequent", "winsorize"], key="imp_method")
                if st.button("Применить", key="imp_apply"):
                    if not cols_imp:
                        st.warning("Выберите хотя бы одну колонку.")
                    else:
                        tmp = df.copy()
                        # применяем выбранный метод
                        for cc in cols_imp:
                            if method_imp == "median":
                                try:
                                    med = tmp[cc].median()
                                    tmp[cc] = tmp[cc].fillna(med)
                                except Exception:
                                    pass
                            elif method_imp == "most_frequent":
                                try:
                                    m = tmp[cc].mode()
                                    if not m.empty:
                                        tmp[cc] = tmp[cc].fillna(m.iloc[0])
                                except Exception:
                                    pass
                            else:  # winsorize
                                try:
                                    tmp[cc] = winsorize_series(tmp[cc], 0.01, 0.99)
                                except Exception:
                                    pass
                        st.dataframe(tmp.head(50))
                        st.download_button(
                            "📥 Скачать (CSV)",
                            data=tmp.to_csv(index=False).encode("utf-8"),
                            file_name="cleaned_data.csv",
                            mime="text/csv",
                        )

    # ---- FEATURES
    with tab_feats:
        st.header("5️⃣ Генератор фичей")
        with st.expander("Фичи для временных рядов", expanded=False):
            dcol = st.selectbox("Дата-колонка", ["<нет>"] + analysis_cols, key="fe_dt")
            if dcol != "<нет>" and numeric_cols:
                vcol = st.selectbox("Метрика", numeric_cols, key="fe_val")
                win = st.slider("Окно для rolling", 3, 60, 7, 1, key="fe_win")
                spike = st.slider("Порог всплеска |pct_change|, %", 5.0, 300.0, 50.0, 5.0, key="fe_spike")
                if st.button("Сгенерировать фичи", key="fe_btn"):
                    feats, g = generate_ts_features(df, dcol, vcol, window=win, spike_thresh_pct=spike)
                    st.json(g)
                    st.dataframe(feats.head(20))
                    st.download_button("📥 Скачать фичи (CSV)", data=feats.to_csv(index=False).encode("utf-8"),
                                       file_name=f"ts_features_{vcol}.csv", mime="text/csv")
            elif dcol != "<нет>":
                st.info("Нужны числовые столбцы.")

        st.markdown("---")
        st.subheader("Табличные фичи")
        with st.expander("Биннинг числовых признаков", expanded=False):
            if numeric_cols:
                bcol = st.selectbox("Колонка", numeric_cols, key="bin_col")
                method = st.radio("Метод", ["qcut (по квантилям)", "cut (равная ширина)"], horizontal=True, key="bin_m")
                n_bins = st.slider("Количество бинов", 3, 10, 5, 1, key="bin_n")
                new_name = st.text_input("Имя новой колонки", value=f"{bcol}_bin_{n_bins}", key="bin_name")
                if st.button("Сгенерировать бины", key="bin_btn"):
                    try:
                        binned = (pd.qcut(df[bcol], q=n_bins, duplicates="drop") if method.startswith("qcut")
                                  else pd.cut(df[bcol], bins=n_bins))
                        tmp = df.copy(); tmp[new_name] = binned
                        st.dataframe(tmp[[bcol, new_name]].head(30))
                        st.download_button("📥 Скачать (CSV)", data=tmp.to_csv(index=False).encode("utf-8"),
                                           file_name=f"binned_{new_name}.csv", mime="text/csv")
                    except Exception as e:
                        st.error(f"Биннинг не выполнен: {e}")
            else:
                st.info("Нет числовых колонок.")

        with st.expander("Флаги по порогам", expanded=False):
            if numeric_cols:
                tcol = st.selectbox("Колонка", numeric_cols, key="thr_col")
                ttype = st.radio("Тип", ["> X", "< X"], horizontal=True, key="thr_type")
                tval = st.number_input("Порог X", value=float(df[tcol].median()) if len(df[tcol].dropna()) else 0.0, key="thr_val")
                fname = st.text_input("Имя флага", value=f"{tcol}_flag", key="thr_name")
                if st.button("Создать флаг", key="thr_btn"):
                    tmp = df.copy()
                    tmp[fname] = (tmp[tcol] > tval).astype(int) if ttype == "> X" else (tmp[tcol] < tval).astype(int)
                    st.dataframe(tmp[[tcol, fname]].head(30))
                    st.download_button("📥 Скачать (CSV)", data=tmp.to_csv(index=False).encode("utf-8"),
                                       file_name=f"flag_{fname}.csv", mime="text/csv")
            else:
                st.info("Нет числовых колонок.")

        with st.expander("Логарифмирование и Z-score", expanded=False):
            num_multi = st.multiselect("Колонки", numeric_cols, key="logz_cols")
            if num_multi and st.button("Сгенерировать", key="logz_btn"):
                tmp = df.copy()
                for c in num_multi:
                    tmp[f"{c}_log1p"] = np.log1p(tmp[c])
                    std = tmp[c].std(ddof=0)
                    tmp[f"{c}_z"] = (tmp[c] - tmp[c].mean()) / std if std and not np.isnan(std) else np.nan
                st.dataframe(tmp.head(20))
                st.download_button("📥 Скачать (CSV)", data=tmp.to_csv(index=False).encode("utf-8"),
                                   file_name="transformed_features.csv", mime="text/csv")

    # ---- CATEGORICAL
    with tab_cats:
        st.header("6️⃣ Категориальные признаки")
        if not categorical_cols:
            st.info("Нет категориальных колонок.")
        else:
            cat = st.selectbox("Колонка", categorical_cols)
            with st.expander("Распределение категорий", expanded=True):
                vc = df[cat].astype(str).value_counts(dropna=False)
                total = int(vc.sum())
                freq_df = vc.rename("count").to_frame(); freq_df["share_%"] = freq_df["count"] / total * 100.0
                st.dataframe(freq_df.head(50))
                top = freq_df.head(20).reset_index().rename(columns={"index": cat})
                fig = px.bar(top, x=cat, y="count", title=f"Топ-20 категорий в {cat}"); fig.update_xaxes(tickangle=60)
                st.plotly_chart(fig, use_container_width=True)
            with st.expander("Метрика по категориям", expanded=False):
                if numeric_cols:
                    m = st.selectbox("Метрика", numeric_cols, key="cat_metric")
                    agg = st.selectbox("Агрегат", ["mean", "median", "sum", "count", "std"], key="cat_agg")
                    grouped = df.groupby(cat)[m].agg(agg).sort_values(ascending=False).rename(agg)
                    st.dataframe(grouped.head(50))
                    top_g = grouped.head(20).reset_index()
                    fig2 = px.bar(top_g, x=cat, y=agg, title=f"{m} по {cat}"); fig2.update_xaxes(tickangle=60)
                    st.plotly_chart(fig2, use_container_width=True)
                else:
                    st.info("Нет числовых колонок.")

    # ---- GROUPS / PIVOT
    with tab_groups:
        st.header("7️⃣ Сравнение групп и сводные таблицы")
        if categorical_cols and numeric_cols:
            with st.expander("Сравнение групп (groupby)", expanded=True):
                gcol = st.selectbox("Группа", categorical_cols, key="grp_col")
                mcol = st.selectbox("Метрика", numeric_cols, key="grp_metric")
                agg = st.selectbox("Агрегат", ["mean", "median", "sum", "count", "std"], key="grp_agg")
                grouped = df.groupby(gcol)[mcol].agg(agg).sort_values(ascending=False).rename(agg)
                st.dataframe(grouped.head(100))
                top_g = grouped.head(20).reset_index()
                st.plotly_chart(px.bar(top_g, x=gcol, y=agg, title=f"Топ-20 {gcol} по {mcol}"), use_container_width=True)
        else:
            st.info("Для сравнения групп нужны и категориальные, и числовые.")

        st.markdown("---")
        st.subheader("📊 Pivot / сводная таблица")
        if categorical_cols and numeric_cols:
            row_col = st.selectbox("Строки (rows)", categorical_cols, key="pivot_row")
            col_col = st.selectbox("Столбцы (columns, опционально)", ["<нет>"] + categorical_cols, key="pivot_col")
            val_col = st.selectbox("Значение (value)", numeric_cols, key="pivot_val")
            agg_pivot = st.selectbox("Агрегация", ["mean", "sum", "count", "median", "std"], key="pivot_agg")

            with st.expander("Доп. фильтры (по любым колонкам)", expanded=False):
                extra_cols = st.multiselect("Колонки для доп. фильтрации", options=df.columns.tolist(), key="pivot_extra_cols")
                extra_filters = {}
                for c in extra_cols:
                    s = df[c]
                    if np.issubdtype(s.dtype, np.number):
                        mn, mx = float(s.min()), float(s.max())
                        if not (np.isnan(mn) or np.isnan(mx)):
                            extra_filters[c] = ("range", st.slider(f"Диапазон для {c}", mn, mx, (mn, mx), key=f"pf_rng_{c}"))
                    else:
                        vals = sorted(s.dropna().astype(str).unique())
                        chosen = st.multiselect(f"Значения для {c}", options=vals, key=f"pf_vals_{c}")
                        if chosen: extra_filters[c] = ("values", chosen)

            with st.expander("Фильтры по rows/columns", expanded=False):
                row_vals = st.multiselect(f"Фильтр по значениям {row_col}", options=sorted(df[row_col].dropna().astype(str).unique()), key="pf_row_vals")
                col_vals = st.multiselect(f"Фильтр по значениям {col_col}", options=sorted(df[col_col].dropna().astype(str).unique()), key="pf_col_vals") if col_col != "<нет>" else []

            chart_type = st.selectbox("Тип графика", [
                "Без графика",
                "Heatmap",
                "Bar (rows, вертикальный)",
                "Bar (rows, горизонтальный)",
                "Bar (columns, вертикальный)",
                "Bar (columns, горизонтальный)",
                "Line (rows)",
                "Line (columns)",
                "Stacked bar (rows × columns)",
                "Stacked bar (columns × rows)",
                "Treemap (rows)",
            ], key="pivot_chart")

            if st.button("Построить pivot"):
                d = df.copy()
                if row_vals: d = d[d[row_col].astype(str).isin(row_vals)]
                if col_col != "<нет>" and col_vals: d = d[d[col_col].astype(str).isin(col_vals)]
                for c, (kind, val) in extra_filters.items():
                    if kind == "range":
                        lo, hi = val; d = d[(d[c] >= lo) & (d[c] <= hi)]
                    else:
                        d = d[d[c].astype(str).isin(val)]
                columns_arg = None if col_col == "<нет>" else col_col
                pvt = pd.pivot_table(d, index=row_col, columns=columns_arg, values=val_col, aggfunc=agg_pivot)
                st.dataframe(pvt)

                if chart_type != "Без графика" and pvt.size > 0:
                    p = pvt.to_frame("value").reset_index() if isinstance(pvt, pd.Series) else pvt.copy()

                    if chart_type == "Heatmap":
                        if pvt.shape[0] <= 50 and (1 if isinstance(pvt, pd.Series) else pvt.shape[1]) <= 50:
                            st.plotly_chart(px.imshow(pvt if not isinstance(pvt, pd.Series) else pvt.to_frame("value"),
                                                      aspect="auto", color_continuous_scale="Blues",
                                                      title=f"Heatmap: {agg_pivot}({val_col})"),
                                            use_container_width=True)
                        else:
                            st.caption("Слишком большая таблица для heatmap.")

                    elif chart_type in ["Bar (rows, вертикальный)", "Bar (rows, горизонтальный)"]:
                        series_rows = (pvt.sum(axis=1) if not isinstance(pvt, pd.Series) and pvt.shape[1] > 1 else pvt.squeeze())
                        df_bar = series_rows.reset_index(); df_bar.columns = [row_col, "value"]
                        fig = px.bar(df_bar, x=row_col, y="value", title=f"{agg_pivot}({val_col}) по {row_col}")
                        if chart_type.endswith("горизонтальный"):
                            fig.update_traces(orientation="h"); fig.update_yaxes(categoryorder="total ascending")
                        fig.update_xaxes(tickangle=60); st.plotly_chart(fig, use_container_width=True)

                    elif chart_type in ["Bar (columns, вертикальный)", "Bar (columns, горизонтальный)"] and columns_arg is not None:
                        series_cols = pvt.sum(axis=0)
                        df_bar = series_cols.reset_index(); df_bar.columns = [col_col, "value"]
                        fig = px.bar(df_bar, x=col_col, y="value", title=f"{agg_pivot}({val_col}) по {col_col}")
                        if chart_type.endswith("горизонтальный"):
                            fig.update_traces(orientation="h"); fig.update_yaxes(categoryorder="total ascending")
                        fig.update_xaxes(tickangle=60); st.plotly_chart(fig, use_container_width=True)

                    elif chart_type == "Line (rows)":
                        series_rows = pvt.sum(axis=1)
                        df_line = series_rows.reset_index(); df_line.columns = [row_col, "value"]
                        st.plotly_chart(px.line(df_line, x=row_col, y="value", markers=True,
                                                title=f"Линейный график по {row_col}"), use_container_width=True)

                    elif chart_type == "Line (columns)" and columns_arg is not None:
                        series_cols = pvt.sum(axis=0)
                        df_line = series_cols.reset_index(); df_line.columns = [col_col, "value"]
                        st.plotly_chart(px.line(df_line, x=col_col, y="value", markers=True,
                                                title=f"Линейный график по {col_col}"), use_container_width=True)

                    elif chart_type == "Stacked bar (rows × columns)" and columns_arg is not None:
                        df_melt = pvt.reset_index().melt(id_vars=[row_col], var_name=col_col, value_name="value")
                        st.plotly_chart(px.bar(df_melt, x=row_col, y="value", color=col_col,
                                               title=f"Stacked: {row_col} × {col_col}"), use_container_width=True)

                    elif chart_type == "Stacked bar (columns × rows)" and columns_arg is not None:
                        df_melt = pvt.reset_index().melt(id_vars=[row_col], var_name=col_col, value_name="value")
                        st.plotly_chart(px.bar(df_melt, x=col_col, y="value", color=row_col,
                                               title=f"Stacked: {col_col} × {row_col}"), use_container_width=True)

                    elif chart_type == "Treemap (rows)":
                        series_rows = (pvt.sum(axis=1) if not isinstance(pvt, pd.Series) and pvt.shape[1] > 1 else pvt.squeeze())
                        df_tree = series_rows.reset_index(); df_tree.columns = [row_col, "value"]
                        st.plotly_chart(px.treemap(df_tree, path=[row_col], values="value", title=f"Treemap по {row_col}"),
                                        use_container_width=True)
        else:
            st.info("Для сводных таблиц нужна хотя бы 1 категориальная и 1 числовая колонка.")

    # ---- A/B
    with tab_ab:
        st.header("8️⃣ A/B тест")
        if categorical_cols and numeric_cols:
            with st.expander("Настройки", expanded=True):
                gcol = st.selectbox("Колонка групп", categorical_cols, key="ab_col")
                levels = sorted(df[gcol].dropna().astype(str).unique().tolist())
                if len(levels) < 2:
                    st.info("Нужно ≥2 разных значения."); gA=gB=None
                else:
                    gA = st.selectbox("Группа A", levels, key="ab_A")
                    gB = st.selectbox("Группа B", [v for v in levels if v != gA], key="ab_B") if len(levels)>1 else None
                m = st.selectbox("Метрика", numeric_cols, key="ab_m")
                test_kind = st.selectbox("Тип теста", ["t-test (Welch)", "t-test (equal var)", "Mann–Whitney U", "z-test (proportions)"])
                alpha = st.number_input(
                    "α",
                    min_value=0.001,
                    max_value=0.2,
                    value=0.05,
                    step=0.005,
                    help=(
                        "Уровень значимости α: вероятность ошибочного отклонения нулевой гипотезы. "
                        "Стандартное значение 0.05."
                    ),
                )

            if st.button("Запустить тест"):
                if not gA or not gB:
                    st.warning("Выберите 2 группы.")
                else:
                    x = df.loc[df[gcol].astype(str)==gA, m].dropna()
                    y = df.loc[df[gcol].astype(str)==gB, m].dropna()
                    if len(x)<2 or len(y)<2:
                        st.info("Слишком мало наблюдений.")
                    else:
                        mean_a, mean_b, diff = float(x.mean()), float(y.mean()), float(y.mean() - x.mean())
                        d_val = cohen_d(x, y)
                        if test_kind == "Mann–Whitney U":
                            # Непараметрический тест для сравнения распределений
                            u, p = stats.mannwhitneyu(x, y, alternative="two-sided")
                            st.subheader("Mann–Whitney U")
                            st.write({"u_stat": float(u), "p_value": float(p)})
                        elif test_kind == "z-test (proportions)":
                            # Тест на разницу долей (предполагаем бинарную метрику, 0/1)
                            try:
                                # Доли успехов
                                p1 = x.mean() if len(x) else np.nan
                                p2 = y.mean() if len(y) else np.nan
                                p_hat = (x.sum() + y.sum()) / (len(x) + len(y)) if (len(x) + len(y)) > 0 else np.nan
                                se = math.sqrt(p_hat * (1 - p_hat) * (1 / len(x) + 1 / len(y))) if (len(x) > 0 and len(y) > 0) else np.nan
                                z_stat = (p2 - p1) / se if se and not np.isnan(se) else np.nan
                                p = 2 * (1 - stats.norm.cdf(abs(z_stat))) if not np.isnan(z_stat) else np.nan
                                st.subheader("z-test (proportions)")
                                st.write({"z_stat": float(z_stat) if not np.isnan(z_stat) else None, "p_value": float(p) if not np.isnan(p) else None})
                            except Exception as _:
                                p = np.nan
                        else:
                            # t-тесты (Welch или с равными дисперсиями)
                            equal_var = (test_kind == "t-test (equal var)")
                            t, p = stats.ttest_ind(x, y, equal_var=equal_var)
                            st.subheader("t-test")
                            st.write({"t_stat": float(t), "p_value": float(p), "cohen_d": float(d_val) if not np.isnan(d_val) else None})
                        c1, c2, c3, c4 = st.columns(4)
                        c1.metric(f"Mean A ({gA})", f"{mean_a:.3g}"); c2.metric(f"Mean B ({gB})", f"{mean_b:.3g}")
                        c3.metric("Δ (B-A)", f"{diff:.3g}"); c4.metric("p-value", f"{p:.4f}")
                        local_ab_summary = business_summary_for_ab(gA, gB, mean_a, mean_b, diff, p, alpha, d_val)
                        render_ai_block(
                            local_ab_summary,
                            "🤖 AI-бизнес-резюме по A/B-тесту",
                            f"ab_main_{gA}_{gB}",
                            extra_prompt=(
                                "Сделай бизнес-резюме по результатам этого A/B-теста. "
                                "Опиши, есть ли статистически значимые отличия, каков размер эффекта "
                                "и какие действия можно рекомендовать продакт-менеджеру."
                            ),
                        )
        else:
            st.info("Нужны и категориальные, и числовые колонки.")

    # ---- SCENARIOS
    with tab_scenarios:
        st.header("9️⃣ Сценарии / шаблоны")
        with st.expander("Быстрый EDA по выбранным метрикам", expanded=True):
            if numeric_cols:
                eda_cols = st.multiselect("Метрики", numeric_cols, default=numeric_cols[:min(3, len(numeric_cols))])
                if eda_cols:
                    target_metric = st.selectbox("Основная метрика для бизнес-резюме", eda_cols, index=0)
                else:
                    target_metric = None

                if st.button("Запустить быстрый EDA"):
                    if not eda_cols:
                        st.warning("Выберите метрики.")
                    else:
                        # 1) Базовые статистики
                        stats_df = pd.DataFrame({c: describe_basic_stats(df[c]) for c in eda_cols}).T
                        with st.expander("1) Базовые статистики", expanded=True):
                            st.dataframe(stats_df)

                        # 2) Корреляции
                        corr_matrix = None
                        if len(eda_cols) >= 2:
                            # Кэшируем расчёт корреляций для ускорения
                            corr_matrix = compute_corr_matrix_cached(df, tuple(eda_cols))
                            with st.expander("2) Корреляции", expanded=False):
                                st.plotly_chart(
                                    px.imshow(
                                        corr_matrix,
                                        text_auto=False,
                                        color_continuous_scale="RdBu",
                                        zmin=-1,
                                        zmax=1,
                                        title="Корреляции (быстрый EDA)",
                                    ),
                                    use_container_width=True,
                                )

                        # 3) Качество данных
                        dq = compute_data_quality_table(df[eda_cols])
                        with st.expander("3) Качество данных", expanded=False):
                            st.dataframe(dq)

                        # 4) Авто-сводка по EDA (локальная + AI по кнопке)
                        local_summary = auto_eda_summary(df, stats_df, corr_matrix, dq, eda_cols)
                        with st.expander("4) Авто-сводка по EDA", expanded=True):
                            render_ai_block(
                                local_summary,
                                "🤖 AI-сводка по EDA",
                                "scenario_eda_global",
                                extra_prompt=(
                                    "На основе приведённых статистик, корреляций и качества данных "
                                    "сделай краткую аналитическую сводку для продакт-менеджера. "
                                    "Опиши основные риски, аномалии и связи между метриками."
                                ),
                            )

                        # 5) Бизнес-резюме по основной метрике
                        if target_metric is not None:
                            metric_stats = stats_df.loc[target_metric].to_dict()
                            local_metric_summary = business_summary_for_series(
                                target_metric,
                                metric_stats,
                                norm_res=None,
                                n_outliers=0,
                            )
                            with st.expander("5) 🧠 Бизнес-резюме по метрике", expanded=True):
                                render_ai_block(
                                    local_metric_summary,
                                    "🤖 AI-резюме по основной метрике",
                                    f"scenario_metric_{target_metric}",
                                    extra_prompt=(
                                        f"Сделай бизнес-резюме по метрике '{target_metric}' на основе статистик и контекста. "
                                        "Опиши, как её поведение может влиять на продукт/бизнес, какие гипотезы и действия стоит проверить."
                                    ),
                                )

                        # Для HTML-отчёта берём AI-сводку, если она уже посчитана, иначе локальную
                        summary_for_report = local_summary
                        if "ai_summaries" in st.session_state:
                            summary_for_report = st.session_state["ai_summaries"].get("scenario_eda_global", local_summary)

                        report_bytes = build_auto_eda_html(
                            df, eda_cols, stats_df, corr_matrix, dq, summary_for_report
                        )
                        st.download_button(
                            "📥 Скачать авто-EDA отчёт (HTML)",
                            data=report_bytes,
                            file_name="auto_eda_report.html",
                            mime="text/html",
                        )
            else:
                st.info("Нет числовых колонок для EDA.")
        with st.expander("Быстрый A/B (две крупнейшие группы)", expanded=False):
            if categorical_cols and numeric_cols:
                gcol = st.selectbox("Колонка групп"
, categorical_cols, key="sc_g")
                m = st.selectbox("Метрика", numeric_cols, key="sc_m")
                levels = df[gcol].astype(str).value_counts().index.tolist()
                if len(levels) >= 2:
                    gA, gB = levels[0], levels[1]
                    st.caption(f"A={gA}, B={gB}")
                    if st.button("Запустить быстрый A/B"):
                        x = df.loc[df[gcol].astype(str)==gA, m].dropna()
                        y = df.loc[df[gcol].astype(str)==gB, m].dropna()
                        if len(x)>=2 and len(y)>=2:
                            mean_a, mean_b, diff = float(x.mean()), float(y.mean()), float(y.mean()-x.mean())
                            t, p = stats.ttest_ind(x, y, equal_var=False); d = cohen_d(x, y)
                            c1, c2, c3, c4 = st.columns(4)
                            c1.metric(f"Mean A ({gA})", f"{mean_a:.3g}"); c2.metric(f"Mean B ({gB})", f"{mean_b:.3g}")
                            c3.metric("Δ (B-A)", f"{diff:.3g}"); c4.metric("p-value", f"{p:.4f}")
                            st.write({"t_stat": float(t), "p_value": float(p), "cohen_d": float(d) if not np.isnan(d) else None})

                            # Локальное бизнес-резюме по быстрому A/B + AI по кнопке
                            local_ab_summary = business_summary_for_ab(gA, gB, mean_a, mean_b, diff, p, 0.05, d)
                            render_ai_block(
                                local_ab_summary,
                                "🤖 AI-резюме по быстрому A/B",
                                f"ab_quick_{gcol}_{m}",
                                extra_prompt=(
                                    "Сделай бизнес-резюме по результатам этого A/B-теста. "
                                    "Опиши, есть ли статистически значимые отличия, каков размер эффекта "
                                    "и какие действия можно рекомендовать продакт-менеджеру."
                                ),
                            )
                        else:
                            st.info("Слишком мало наблюдений.")
                else:
                    st.info("Недостаточно групп.")
            else:
                st.info("Нужны и категориальные, и числовые колонки.")

        # --- Новые сценарии / шаблоны ---
        with st.expander("Диагностика каннибализации", expanded=False):
            """Анализ возможной каннибализации между двумя категориями. Выберите категориальную колонку,
            метрику и колонку даты; далее выберите две категории, для которых сравниваются временные
            ряды. Отрицательная корреляция может свидетельствовать о каннибализации."""
            if categorical_cols and numeric_cols:
                cannib_cat = st.selectbox(
                    "Категориальная колонка",
                    categorical_cols,
                    key="cannib_cat",
                )
                cannib_metric = st.selectbox(
                    "Метрика",
                    numeric_cols,
                    key="cannib_metric",
                )
                cannib_date = st.selectbox(
                    "Колонка даты/времени",
                    ["<нет>"] + analysis_cols,
                    key="cannib_date",
                )
                if cannib_date != "<нет>":
                    categories = df[cannib_cat].dropna().astype(str).unique().tolist()
                    if len(categories) >= 2:
                        cat1 = st.selectbox("Категория 1", categories, key="cannib_cat1")
                        cat2_opts = [c for c in categories if c != cat1]
                        cat2 = st.selectbox("Категория 2", cat2_opts, key="cannib_cat2") if cat2_opts else None
                        if cat2:
                            if st.button("Запустить анализ каннибализации", key="btn_cannib"):
                                try:
                                    d = df[[cannib_date, cannib_cat, cannib_metric]].dropna().copy()
                                    d[cannib_date] = pd.to_datetime(d[cannib_date])
                                    s1 = (
                                        d[d[cannib_cat] == cat1]
                                        .groupby(cannib_date)[cannib_metric]
                                        .sum()
                                        .rename(cat1)
                                    )
                                    s2 = (
                                        d[d[cannib_cat] == cat2]
                                        .groupby(cannib_date)[cannib_metric]
                                        .sum()
                                        .rename(cat2)
                                    )
                                    joined = pd.concat([s1, s2], axis=1).dropna()
                                    if len(joined) >= 2:
                                        corr_val = joined[cat1].corr(joined[cat2])
                                        st.line_chart(joined)
                                        st.metric("Корреляция (r)", f"{corr_val:.3g}")
                                        st.caption(
                                            "Отрицательная корреляция между временными рядами может указывать на каннибализацию."
                                        )
                                    else:
                                        st.info("Недостаточно данных для анализа.")
                                except Exception as e:
                                    st.error(f"Ошибка анализа: {e}")
                    else:
                        st.info("Недостаточно уникальных категорий.")
                else:
                    st.info("Выберите колонку даты/времени для анализа каннибализации.")
            else:
                st.info("Нужна хотя бы одна категориальная и одна числовая колонка.")

        with st.expander("Выявление сезонности (ACF/STL)", expanded=False):
            """Определение сезонности с помощью автокорреляционной функции и STL-разложения. Выберите дату,
            метрику и максимальный лаг. Программа покажет ACF и предложит возможный сезонный период."""
            if numeric_cols:
                seas_date = st.selectbox(
                    "Колонка даты",
                    ["<нет>"] + analysis_cols,
                    key="seas_date",
                )
                seas_value = st.selectbox("Метрика", numeric_cols, key="seas_value")
                seas_maxlag = st.slider(
                    "Максимальный лаг для ACF",
                    min_value=5,
                    max_value=100,
                    value=30,
                    step=1,
                    key="seas_maxlag",
                )
                if st.button("Рассчитать сезонность", key="btn_seas"):
                    if seas_date == "<нет>":
                        st.warning("Выберите колонку даты.")
                    else:
                        try:
                            d = df[[seas_date, seas_value]].dropna().copy()
                            d[seas_date] = pd.to_datetime(d[seas_date])
                            d = d.sort_values(seas_date)
                            lags, acf_vals = compute_acf(d[seas_value], seas_maxlag)
                            fig = go.Figure()
                            fig.add_trace(go.Bar(x=lags, y=acf_vals, name="ACF"))
                            fig.update_layout(
                                title="Автокорреляционная функция",
                                xaxis_title="Lag",
                                yaxis_title="ACF",
                                height=400,
                            )
                            st.plotly_chart(fig, use_container_width=True)
                            # Определяем вероятный сезонный период (исключаем лаг 0)
                            if len(acf_vals) > 1:
                                best_idx = int(np.argmax(np.abs(acf_vals[1:])) + 1)
                                best_lag = int(lags[best_idx])
                                st.write(f"Возможный сезонный период: {best_lag}")
                            # STL-разложение по выбору пользователя
                            if STL is not None and st.checkbox("Показать STL-разложение", key="seas_stl"):
                                try:
                                    comp = STL(d[seas_value], period=best_lag if len(acf_vals) > 1 else seas_maxlag).fit()
                                    comp_df = pd.DataFrame(
                                        {
                                            "date": d[seas_date].values,
                                            "observed": comp.observed,
                                            "trend": comp.trend,
                                            "seasonal": comp.seasonal,
                                            "resid": comp.resid,
                                        }
                                    )
                                    fig_stl = make_stl_figure(comp_df)
                                    st.plotly_chart(fig_stl, use_container_width=True)
                                except Exception as e:
                                    st.error(f"Ошибка STL: {e}")
                        except Exception as e:
                            st.error(f"Ошибка ACF: {e}")
            else:
                st.info("Нужна хотя бы одна числовая колонка.")

        with st.expander("Влияние акций / событий", expanded=False):
            """Оценка влияния маркетинговой акции или события: сравнение среднего значения метрики до и после указанной даты."
            """
            if numeric_cols:
                event_date_col = st.selectbox(
                    "Колонка даты",
                    ["<нет>"] + analysis_cols,
                    key="event_date_col",
                )
                event_value = st.selectbox(
                    "Метрика",
                    numeric_cols,
                    key="event_value_col",
                )
                event_dt = st.date_input(
                    "Дата события (разделение до/после)",
                    key="event_dt",
                )
                if st.button("Оценить влияние", key="btn_event"):
                    if event_date_col == "<нет>":
                        st.warning("Выберите колонку даты.")
                    else:
                        try:
                            d = df[[event_date_col, event_value]].dropna().copy()
                            d[event_date_col] = pd.to_datetime(d[event_date_col])
                            before = d[d[event_date_col] < pd.to_datetime(event_dt)][event_value]
                            after = d[d[event_date_col] >= pd.to_datetime(event_dt)][event_value]
                            if len(before) >= 2 and len(after) >= 2:
                                mean_before = float(before.mean())
                                mean_after = float(after.mean())
                                diff = mean_after - mean_before
                                t_stat, p_val = stats.ttest_ind(after, before, equal_var=False)
                                c1, c2, c3, c4 = st.columns(4)
                                c1.metric("Среднее до", f"{mean_before:.3g}")
                                c2.metric("Среднее после", f"{mean_after:.3g}")
                                c3.metric("Δ (после - до)", f"{diff:.3g}")
                                c4.metric("p-value", f"{p_val:.4f}")
                                if p_val < 0.05:
                                    st.success("Изменение статистически значимо.")
                                else:
                                    st.info("Статистически значимых изменений не выявлено.")
                            else:
                                st.info("Недостаточно наблюдений до или после выбранной даты.")
                        except Exception as e:
                            st.error(f"Ошибка анализа: {e}")
            else:
                st.info("Нет числовых колонок для анализа.")

    # ---- DICTIONARY
    with tab_dict:
        st.header("🔟 Справочник статистических методов")
        key = st.selectbox("Метод", list(METHOD_INFO.keys()), format_func=lambda k: METHOD_INFO[k]["name"])
        info = METHOD_INFO[key]
        st.subheader(info["name"]); st.markdown(f"**Описание:** {info['description']}"); st.markdown(f"**Когда использовать:** {info['when']}")


    # ---- REPORT
    with tab_report:
        st.header("1️⃣1️⃣ Отчёты и экспорт")

        if df is None or df.empty:
            st.info("Нет данных для отчёта. Загрузите и подготовьте датасет на шагах 1–2.")
        else:
            st.subheader("🔧 Конфиг отчёта")

            # Какие числовые метрики включать в EDA-часть
            num_cols = df.select_dtypes(include="number").columns.tolist()
            report_num_cols = st.multiselect(
                "Числовые метрики для отчёта (EDA, корреляции, AI-сводка)",
                num_cols,
                default=num_cols,
                key="report_num_cols",
            )

            if not report_num_cols:
                st.info("Выберите хотя бы одну числовую метрику для отчёта.")
            else:
                cols = report_num_cols
                stats_df = pd.DataFrame({c: describe_basic_stats(df[c]) for c in cols}).T

                corr_matrix = None
                if len(cols) >= 2:
                    try:
                        # Используем кэшированную функцию корреляции для ускорения
                        corr_matrix = compute_corr_matrix_cached(df, tuple(cols))
                    except Exception:
                        corr_matrix = None

                dq = compute_data_quality_table(df[cols])

                # Текстовая авто-сводка (локальная) + опциональная AI-версия по кнопке
                summary_text = auto_eda_summary(df, stats_df, corr_matrix, dq, cols)
                st.subheader("📄 Локальная сводка для отчёта")
                st.markdown(f"<div class='business-summary'>{summary_text}</div>", unsafe_allow_html=True)

                with st.expander("🤖 AI-сводка для отчёта", expanded=False):
                    render_ai_block(
                        summary_text,
                        "🤖 Сгенерировать общую AI-сводку для отчёта",
                        "report_global",
                        extra_prompt=(
                            "Сделай краткое бизнес-резюме для презентации менеджменту. "
                            "Сфокусируйся на ключевых рисках, возможностях роста и гипотезах для A/B-тестов."
                        ),
                    )

                # Выбираем текст, который пойдёт в отчёты: AI, если есть, иначе локальный
                summary_for_report = summary_text
                if "ai_summaries" in st.session_state:
                    summary_for_report = st.session_state["ai_summaries"].get("report_global", summary_text)

                # --- Формирование и скачивание отчётов ---
                # HTML‑отчёт (авто‑EDA)
                html_bytes = build_auto_eda_html(
                    df,
                    cols,
                    stats_df,
                    corr_matrix,
                    dq,
                    summary_text=summary_for_report,
                )
                st.download_button(
                    "📥 Скачать HTML-отчёт (EDA + AI-сводка)",
                    data=html_bytes,
                    file_name="auto_eda_ai_report.html",
                    mime="text/html",
                )

                # Excel‑отчёт (Data + EDA + DQ + AI). Создаётся по нажатию кнопки.
                report_xlsx: BytesIO | None = None
                report_pptx: BytesIO | None = None
                if st.button("Сформировать отчёты (Excel + PPTX)", key="btn_build_excel_report"):
                    with st.spinner("Генерируем отчёты..."):
                        text_blocks = {
                            "EDA summary (rule-based)": summary_text,
                            "Global AI summary": summary_for_report,
                        }
                        report_xlsx = build_excel_report(
                            df,
                            stats_df=stats_df,
                            corr_matrix=corr_matrix,
                            dq=dq,
                            text_blocks=text_blocks,
                        )
                        # Строим PPTX-отчёт
                        report_pptx = build_pptx_report(
                            df,
                            stats_df=stats_df,
                            corr_matrix=corr_matrix,
                            dq=dq,
                            summary_text=summary_for_report,
                        )
                    if report_xlsx:
                        st.download_button(
                            "📥 Скачать Excel-отчёт",
                            data=report_xlsx,
                            file_name="stats_lab_report_v5_6.xlsx",
                            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                        )
                    if report_pptx and report_pptx.getbuffer().nbytes > 0:
                        st.download_button(
                            "📥 Скачать PPTX-отчёт",
                            data=report_pptx,
                            file_name="stats_lab_report_v5_6.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        )

                # One‑click отчёт: zip HTML, Excel и PPTX, если доступны
                if html_bytes and report_xlsx:
                    try:
                        import zipfile
                        bundle = BytesIO()
                        with zipfile.ZipFile(bundle, mode="w") as zf:
                            zf.writestr("report.html", html_bytes)
                            # Excel
                            if hasattr(report_xlsx, "getvalue"):
                                zf.writestr("report.xlsx", report_xlsx.getvalue())
                            else:
                                zf.writestr("report.xlsx", report_xlsx)
                            # PPTX
                            if report_pptx and hasattr(report_pptx, "getvalue") and report_pptx.getbuffer().nbytes > 0:
                                zf.writestr("report.pptx", report_pptx.getvalue())
                        bundle.seek(0)
                        st.download_button(
                            "📥 One‑click отчёт (HTML + Excel + PPTX)",
                            data=bundle.getvalue(),
                            file_name="stats_lab_report_bundle.zip",
                            mime="application/zip",
                        )
                    except Exception:
                        pass

                # Профиль анализа: сохраняем фильтры и выбранные числовые колонки
                profile_dict = {"filters": filter_config, "selected_report_columns": cols}
                profile_bytes = json.dumps(profile_dict, ensure_ascii=False, indent=2).encode("utf-8")
                st.download_button(
                    "📥 Скачать профиль анализа (JSON)",
                    data=profile_bytes,
                    file_name="analysis_profile.json",
                    mime="application/json",
                )

                # Отчёт только с данными (если нужны только DataSheet)
                report = build_excel_report(df, None, None, None, None)
                st.download_button(
                    "📥 Скачать отчёт Excel",
                    data=report,
                    file_name="stats_report_v5_6.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                )


if __name__ == "__main__":
    main()